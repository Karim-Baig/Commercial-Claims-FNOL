"""Single-slide overview of the 8 Claude Code skills in Skills-Agents-Repository."""
import sys, os


def _find_skill_scripts():
    skills_dir = os.path.join(os.path.expanduser("~"), ".claude", "skills")
    for candidate in ["acnpptx", "pptx"]:
        path = os.path.join(skills_dir, candidate, "scripts")
        if os.path.exists(os.path.join(path, "helpers.py")):
            return path
    raise FileNotFoundError("Cannot find acnpptx skill scripts.")


sys.path.insert(0, _find_skill_scripts())
from helpers import *                                   # noqa: F401,F403
from svg_pipeline import add_svg_native as _add_svg_native, cleanup_temp
from pptx import Presentation
from pptx.oxml.ns import qn as _qn
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

prs = Presentation(TEMPLATE_PATH)
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    prs.part.drop_rel(sldId.get(_qn('r:id')))
    del prs.slides._sldIdLst[0]


def add_svg_native(slide, svg_str, x, y, w, h, png_width=800):
    return _add_svg_native(slide, prs, svg_str, x, y, w, h, png_width)


# ── icons (from the skill's SVG icon library) ────────────────────────────────
def svg_icon_layers(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <polygon points="32,8 56,22 32,36 8,22" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>
  <polygon points="32,20 56,34 32,48 8,34" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round" opacity="0.7"/>
  <polygon points="32,32 56,46 32,60 8,46" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round" opacity="0.4"/>
</svg>'''


def svg_icon_shield(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M32 6L8 18v16c0 14 10 22 24 26 14-4 24-12 24-26V18L32 6z" fill="none" stroke="{color}" stroke-width="2.5" stroke-linejoin="round"/>
  <polyline points="22,32 30,40 44,24" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''


def svg_icon_database(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <ellipse cx="32" cy="14" rx="22" ry="8" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M10 14v36c0 4.4 9.8 8 22 8s22-3.6 22-8V14" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M10 26c0 4.4 9.8 8 22 8s22-3.6 22-8" fill="none" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
  <path d="M10 38c0 4.4 9.8 8 22 8s22-3.6 22-8" fill="none" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
</svg>'''


def svg_icon_calendar(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <rect x="8" y="14" width="48" height="42" rx="4" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>
  <line x1="8" y1="26" x2="56" y2="26" stroke="{color}" stroke-width="2"/>
  <line x1="20" y1="8" x2="20" y2="20" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <line x1="44" y1="8" x2="44" y2="20" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <circle cx="20" cy="35" r="2" fill="{color}"/><circle cx="32" cy="35" r="2" fill="{color}"/>
  <circle cx="44" cy="35" r="2" fill="{color}"/><circle cx="20" cy="46" r="2" fill="{color}"/>
  <circle cx="32" cy="46" r="2" fill="{color}"/><circle cx="44" cy="46" r="2" fill="{color}"/>
</svg>'''


def svg_icon_chart_up(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <line x1="8" y1="56" x2="56" y2="56" stroke="#D8D8D8" stroke-width="1.5"/>
  <line x1="8" y1="56" x2="8" y2="8" stroke="#D8D8D8" stroke-width="1.5"/>
  <polyline points="12,48 24,38 36,26 48,12" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
  <circle cx="12" cy="48" r="3" fill="{color}"/><circle cx="24" cy="38" r="3" fill="{color}"/>
  <circle cx="36" cy="26" r="3" fill="{color}"/><circle cx="48" cy="12" r="3" fill="{color}"/>
</svg>'''


def svg_icon_network(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="32" cy="14" r="7" fill="none" stroke="{color}" stroke-width="2.5"/>
  <circle cx="14" cy="50" r="7" fill="none" stroke="{color}" stroke-width="2.5"/>
  <circle cx="50" cy="50" r="7" fill="none" stroke="{color}" stroke-width="2.5"/>
  <circle cx="32" cy="14" r="3" fill="{color}"/><circle cx="14" cy="50" r="3" fill="{color}"/>
  <circle cx="50" cy="50" r="3" fill="{color}"/>
  <line x1="32" y1="21" x2="14" y2="43" stroke="{color}" stroke-width="1.5"/>
  <line x1="32" y1="21" x2="50" y2="43" stroke="{color}" stroke-width="1.5"/>
  <line x1="21" y1="50" x2="43" y2="50" stroke="{color}" stroke-width="1.5"/>
</svg>'''


def svg_icon_code(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <polyline points="20,16 6,32 20,48" fill="none" stroke="{color}" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/>
  <polyline points="44,16 58,32 44,48" fill="none" stroke="{color}" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/>
  <line x1="36" y1="10" x2="28" y2="54" stroke="{color}" stroke-width="2.5" stroke-linecap="round" opacity="0.7"/>
</svg>'''


def svg_icon_document(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M16 4h22l14 14v38a4 4 0 0 1-4 4H16a4 4 0 0 1-4-4V8a4 4 0 0 1 4-4z" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M38 4v14h14" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>
  <line x1="20" y1="28" x2="44" y2="28" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
  <line x1="20" y1="36" x2="44" y2="36" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
  <line x1="20" y1="44" x2="36" y2="44" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
</svg>'''


def svg_icon_gear(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M32 20a12 12 0 1 0 0 24 12 12 0 0 0 0-24zm0 18a6 6 0 1 1 0-12 6 6 0 0 1 0 12z" fill="{color}"/>
  <path d="M34.4 8h-4.8l-1.2 5.2c-1.6.4-3 1.2-4.2 2.2L19.4 13l-3.4 3.4 2.4 4.8c-1 1.2-1.8 2.6-2.2 4.2L11 26.6v4.8l5.2 1.2c.4 1.6 1.2 3 2.2 4.2L16 41.6l3.4 3.4 4.8-2.4c1.2 1 2.6 1.8 4.2 2.2L29.6 50h4.8l1.2-5.2c1.6-.4 3-1.2 4.2-2.2l4.8 2.4 3.4-3.4-2.4-4.8c1-1.2 1.8-2.6 2.2-4.2L53 31.4v-4.8l-5.2-1.2c-.4-1.6-1.2-3-2.2-4.2l2.4-4.8-3.4-3.4-4.8 2.4c-1.2-1-2.6-1.8-4.2-2.2L34.4 8z" fill="none" stroke="{color}" stroke-width="1.5"/>
</svg>'''


def svg_icon_refresh(size=64, color="#7E00FF"):
    # Used for the unit-test tile rather than a checkmark: the circle-with-check reads
    # almost identically to the shield-with-check at 0.44", and the re-run loop is the
    # distinctive thing about that skill anyway.
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M48 16A22 22 0 0 0 12 26" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <path d="M16 48A22 22 0 0 0 52 38" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <polyline points="12,16 12,28 24,28" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
  <polyline points="52,48 52,36 40,36" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''


# ── content ──────────────────────────────────────────────────────────────────
TILES = [
    (svg_icon_layers,   "Architecture builder",
     "A paired CLAUDE.md and Architecture.md, with metrics that verify themselves."),
    (svg_icon_calendar, "POC plan builder",
     "A screen-by-screen, day-by-day build plan. Every claim traced back to an RFP clause."),
    (svg_icon_network,  "RFP UI inventory",
     "Screens, navigation flow, personas and permissions, pulled from documents or a Figma file."),
    (svg_icon_code,     "Figma to prototype",
     "Figma screens become a running, clickable front end with the real colours and type."),
    (svg_icon_shield,   "Compliance check",
     "Predicts what SonarQube, Snyk and Semgrep would flag, before the scan runs. Never edits code."),
    (svg_icon_refresh,  "Unit test generator",
     "Catalogues every test worth writing, ranked by risk, then writes the ones you pick."),
    (svg_icon_chart_up, "Python optimizer",
     "Finds hot loops, N+1 queries and missed vectorisation. Ranked by measured gain."),
    (svg_icon_database, "FNOL field reference",
     "Every intake field and required document for any insurance line, as a formatted workbook."),
]

slide = prs.slides.add_slide(prs.slide_layouts[2])
ph = {p.placeholder_format.idx: p for p in slide.placeholders}

ph[11].text_frame.text = "Skills repository"
for p in ph[11].text_frame.paragraphs:
    p.font.size, p.font.bold, p.font.color.rgb, p.font.name = Pt(14), True, PURPLE, FONT

ph[0].text_frame.text = "Eight skills that turn documents into deliverables"
for p in ph[0].text_frame.paragraphs:
    p.font.size, p.font.bold, p.font.color.rgb, p.font.name = Pt(28), True, BLACK, FONT

ph[10].text_frame.text = ("Each skill reads real inputs and writes a real artefact — a document, a workbook, "
                          "a test suite or a running prototype. None of them stop at a summary.")
for p in ph[10].text_frame.paragraphs:
    p.font.size, p.font.color.rgb, p.font.name = Pt(18), BLACK, FONT

# ── 4 x 2 tile grid ──────────────────────────────────────────────────────────
GAP = 0.16
COLS, ROWS = 4, 2
TW = (CW - (COLS - 1) * GAP) / COLS      # 3.005"
TH = (AH - (ROWS - 1) * GAP) / ROWS      # 2.520"

for i, (icon_fn, title, desc) in enumerate(TILES):
    r, c = divmod(i, COLS)
    x = ML + c * (TW + GAP)
    y = CY + r * (TH + GAP)

    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(TW), Inches(TH))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_LIGHT
    card.line.fill.background()
    card.shadow.inherit = False

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(TW), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_svg_native(slide, icon_fn(64, "#7E00FF"), x + 0.22, y + 0.30, 0.44, 0.44, png_width=256)

    tb = slide.shapes.add_textbox(Inches(x + 0.20), Inches(y + 0.86), Inches(TW - 0.40), Inches(0.42))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.size, run.font.bold, run.font.color.rgb, run.font.name = Pt(16), True, BLACK, FONT

    db = slide.shapes.add_textbox(Inches(x + 0.20), Inches(y + 1.32), Inches(TW - 0.40), Inches(TH - 1.50))
    dtf = db.text_frame
    dtf.word_wrap = True
    dtf.margin_left = dtf.margin_right = dtf.margin_top = dtf.margin_bottom = 0
    dp = dtf.paragraphs[0]
    dp.line_spacing = 1.15
    drun = dp.add_run()
    drun.text = desc
    drun.font.size, drun.font.color.rgb, drun.font.name = Pt(14), TEXT_BODY, FONT

# Single-slide deliverable: this slide IS the cover, so it carries the full logo rather
# than the GT symbol. Placed in the footer band bottom-right, clear of the tile grid
# (which ends at BY) and clear of the slide number.
add_logo_to_cover(slide, x=11.20, y=6.98, w=1.40)
set_footer(slide)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — the POC built with the toolchain
# Every figure below was counted from the codebase (pytest --collect-only,
# CREATE TABLE statements, @router decorators), not estimated.
# ═════════════════════════════════════════════════════════════════════════════
KPIS = [
    ("45", "API endpoints\nacross 12 routers"),
    ("16", "database tables\nSQLite to MySQL"),
    ("144", "automated tests\nall passing"),
    ("7", "front-end screens\nplus 9 sub-panels"),
    ("5", "locales\nRTL-aware"),
]

CAPABILITIES = [
    (svg_icon_shield,   "Scoped claims visibility",
     "Org scope comes only from the JWT. Out-of-scope returns 403, never 404, and is audit-logged."),
    (svg_icon_document, "Document security gating",
     "Three sequential gates on every fetch. The ECM reference never leaves the API."),
    (svg_icon_gear,     "Config-driven UI",
     "Claims-list columns are read from a database table — change them with no rebuild."),
    (svg_icon_refresh,  "Resilient FNOL intake",
     "Durable outbox plus an idempotency key, so a retry cannot create a duplicate claim."),
]

ARTEFACTS = [
    ("CLAUDE.md + Architecture.md", "Architecture builder"),
    ("UI page inventory workbook", "RFP UI inventory"),
    ("Screen-by-screen build plan", "POC plan builder"),
    ("144-test pytest suite", "Unit test generator"),
    ("FNOL field + document spec", "FNOL field reference"),
]

s2 = prs.slides.add_slide(prs.slide_layouts[2])
p2 = {p.placeholder_format.idx: p for p in s2.placeholders}

p2[11].text_frame.text = "Proof point"
for p in p2[11].text_frame.paragraphs:
    p.font.size, p.font.bold, p.font.color.rgb, p.font.name = Pt(14), True, PURPLE, FONT

p2[0].text_frame.text = "The skills produced a running, tested claims POC"
for p in p2[0].text_frame.paragraphs:
    p.font.size, p.font.bold, p.font.color.rgb, p.font.name = Pt(28), True, BLACK, FONT

p2[10].text_frame.text = ("Meridian Claims Copilot — a micro-frontend claims platform on FastAPI and React. "
                          "Every figure below is counted from the codebase, not estimated.")
for p in p2[10].text_frame.paragraphs:
    p.font.size, p.font.color.rgb, p.font.name = Pt(18), BLACK, FONT

# ── band 1: KPI strip ────────────────────────────────────────────────────────
K_H = 1.24
KW = (CW - 4 * GAP) / 5
for i, (num, label) in enumerate(KPIS):
    kx = ML + i * (KW + GAP)
    box = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(kx), Inches(CY), Inches(KW), Inches(K_H))
    box.fill.solid()
    box.fill.fore_color.rgb = LP_BG
    box.line.fill.background()
    box.shadow.inherit = False

    tb = s2.shapes.add_textbox(Inches(kx + 0.14), Inches(CY + 0.10), Inches(KW - 0.28), Inches(K_H - 0.20))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r1 = tf.paragraphs[0].add_run()
    r1.text = num
    r1.font.size, r1.font.bold, r1.font.color.rgb, r1.font.name = Pt(30), True, DEEP_PURPLE, FONT
    para = tf.add_paragraph()
    para.line_spacing = 1.1
    r2 = para.add_run()
    r2.text = label
    r2.font.size, r2.font.color.rgb, r2.font.name = Pt(14), TEXT_BODY, FONT

# ── band 2: capabilities (left) + artefact-to-skill map (right) ──────────────
B2Y = CY + K_H + 0.22
LBL_H = 0.30
BODY_Y = B2Y + LBL_H + 0.06
BODY_H = BY - BODY_Y
LEFT_W = 6.05
RIGHT_X = ML + LEFT_W + 0.20
RIGHT_W = CW - LEFT_W - 0.20


def band_label(x, w, text):
    tb = s2.shapes.add_textbox(Inches(x), Inches(B2Y), Inches(w), Inches(LBL_H))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(14), True, PURPLE, FONT


band_label(ML, LEFT_W, "Four capabilities validated")
band_label(RIGHT_X, RIGHT_W, "What each skill contributed")

CGAP = 0.10
CH = (BODY_H - 3 * CGAP) / 4
for i, (icon_fn, title, desc) in enumerate(CAPABILITIES):
    cy = BODY_Y + i * (CH + CGAP)
    bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(cy), Inches(0.06), Inches(CH))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_svg_native(s2, icon_fn(64, "#7E00FF"), ML + 0.20, cy + (CH - 0.38) / 2, 0.38, 0.38, png_width=256)

    tb = s2.shapes.add_textbox(Inches(ML + 0.72), Inches(cy + 0.02), Inches(LEFT_W - 0.80), Inches(CH - 0.04))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r1 = tf.paragraphs[0].add_run()
    r1.text = title
    r1.font.size, r1.font.bold, r1.font.color.rgb, r1.font.name = Pt(14), True, BLACK, FONT
    para = tf.add_paragraph()
    para.line_spacing = 1.12
    r2 = para.add_run()
    r2.text = desc
    r2.font.size, r2.font.color.rgb, r2.font.name = Pt(14), TEXT_BODY, FONT

n_rows = len(ARTEFACTS) + 1
tbl_shape = s2.shapes.add_table(n_rows, 2, Inches(RIGHT_X), Inches(BODY_Y),
                                Inches(RIGHT_W), Inches(BODY_H))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(RIGHT_W * 0.55)
tbl.columns[1].width = Inches(RIGHT_W - RIGHT_W * 0.55)
row_h = BODY_H / n_rows
for r in tbl.rows:
    r.height = Inches(row_h)

HEADERS = ("Artefact in the POC", "Skill that produces it")
for c, txt in enumerate(HEADERS):
    cell = tbl.cell(0, c)
    cell.text = txt
    cell.fill.solid()
    cell.fill.fore_color.rgb = PURPLE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.10)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.size, run.font.bold, run.font.color.rgb, run.font.name = Pt(14), True, WHITE, FONT

for i, (artefact, skill) in enumerate(ARTEFACTS, start=1):
    for c, txt in enumerate((artefact, skill)):
        cell = tbl.cell(i, c)
        cell.text = txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 else BG_LIGHT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.10)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size, run.font.name = Pt(14), FONT
                run.font.color.rgb = TEXT_BODY
                run.font.bold = (c == 1)

add_gt_to_slide(s2)          # internal slide: GT symbol, never the full logo
set_footer(s2)

# ── save + verify ────────────────────────────────────────────────────────────
output_path = os.path.join(os.getcwd(), "Claude_Code_Skills_Overview.pptx")
prs.save(output_path)
cleanup_temp()

import subprocess
_venv_py = os.path.join(os.path.dirname(SKILL_DIR), ".venv", "Scripts", "python.exe")
_py = _venv_py if os.path.exists(_venv_py) else sys.executable
subprocess.run([_py, os.path.join(SKILL_DIR, "scripts", "verify_pptx.py"), output_path], check=False)
subprocess.run([_py, os.path.join(SKILL_DIR, "scripts", "brand_check.py"), output_path], check=False)
subprocess.run([_py, os.path.join(SKILL_DIR, "scripts", "thumbnail.py"), output_path,
                os.path.join(os.getcwd(), "thumbnails")], check=False)
print("\nSaved:", output_path)
