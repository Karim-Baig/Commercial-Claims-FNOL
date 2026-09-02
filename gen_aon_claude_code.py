import sys, os

def _find_skill_scripts():
    skills_dir = os.path.join(os.path.expanduser("~"), ".claude", "skills")
    for candidate in ["acnpptx", "pptx"]:
        p = os.path.join(skills_dir, candidate, "scripts")
        if os.path.exists(os.path.join(p, "helpers.py")):
            return p
    raise FileNotFoundError("Cannot find acnpptx skill scripts.")

sys.path.insert(0, _find_skill_scripts())
from helpers import *
from svg_pipeline import add_svg_native as _add_svg_native, cleanup_temp
from pptx import Presentation
from pptx.oxml.ns import qn as _qn
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation(TEMPLATE_PATH)
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    rId = sldId.get(_qn('r:id'))
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

def add_svg(slide, svg_str, x, y, w, h, pw=800):
    return _add_svg_native(slide, prs, svg_str, x, y, w, h, pw)

# ── INLINE SVG HELPERS ────────────────────────────────────────────────────

def svg_icon_shield(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<path d="M32 6L8 18v16c0 14 10 22 24 26 14-4 24-12 24-26V18L32 6z" fill="none" stroke="{color}" stroke-width="2.5" stroke-linejoin="round"/>'
            f'<polyline points="22,32 30,40 44,24" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>'
            f'</svg>')

def svg_icon_code(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<polyline points="20,16 6,32 20,48" fill="none" stroke="{color}" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/>'
            f'<polyline points="44,16 58,32 44,48" fill="none" stroke="{color}" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/>'
            f'<line x1="36" y1="10" x2="28" y2="54" stroke="{color}" stroke-width="2.5" stroke-linecap="round" opacity="0.7"/>'
            f'</svg>')

def svg_icon_gear(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<path d="M32 20a12 12 0 1 0 0 24 12 12 0 0 0 0-24zm0 18a6 6 0 1 1 0-12 6 6 0 0 1 0 12z" fill="{color}"/>'
            f'<path d="M34.4 8h-4.8l-1.2 5.2c-1.6.4-3 1.2-4.2 2.2L19.4 13l-3.4 3.4 2.4 4.8c-1 1.2-1.8 2.6-2.2 4.2L11 26.6v4.8l5.2 1.2c.4 1.6 1.2 3 2.2 4.2L16 41.6l3.4 3.4 4.8-2.4c1.2 1 2.6 1.8 4.2 2.2L29.6 50h4.8l1.2-5.2c1.6-.4 3-1.2 4.2-2.2l4.8 2.4 3.4-3.4-2.4-4.8c1-1.2 1.8-2.6 2.2-4.2L53 31.4v-4.8l-5.2-1.2c-.4-1.6-1.2-3-2.2-4.2l2.4-4.8-3.4-3.4-4.8 2.4c-1.2-1-2.6-1.8-4.2-2.2L34.4 8z" fill="none" stroke="{color}" stroke-width="1.5"/>'
            f'</svg>')

def svg_icon_search(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<circle cx="28" cy="28" r="18" fill="none" stroke="{color}" stroke-width="2.5"/>'
            f'<line x1="40.7" y1="40.7" x2="56" y2="56" stroke="{color}" stroke-width="3" stroke-linecap="round"/>'
            f'</svg>')

def svg_icon_document(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<path d="M16 4h22l14 14v38a4 4 0 0 1-4 4H16a4 4 0 0 1-4-4V8a4 4 0 0 1 4-4z" fill="none" stroke="{color}" stroke-width="2"/>'
            f'<path d="M38 4v14h14" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>'
            f'<line x1="20" y1="28" x2="44" y2="28" stroke="{color}" stroke-width="1.5" opacity="0.6"/>'
            f'<line x1="20" y1="36" x2="44" y2="36" stroke="{color}" stroke-width="1.5" opacity="0.6"/>'
            f'<line x1="20" y1="44" x2="36" y2="44" stroke="{color}" stroke-width="1.5" opacity="0.6"/>'
            f'</svg>')

def svg_icon_brain(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<path d="M32 8v48" stroke="{color}" stroke-width="1.5" stroke-linecap="round" opacity="0.4"/>'
            f'<path d="M32 12c-4-4-10-4-13 0s-4 10 0 14" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>'
            f'<path d="M19 26c-5 2-8 8-5 13" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>'
            f'<path d="M14 39c-2 5 1 11 7 13h11" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>'
            f'<path d="M32 12c4-4 10-4 13 0s4 10 0 14" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>'
            f'<path d="M45 26c5 2 8 8 5 13" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>'
            f'<path d="M50 39c2 5-1 11-7 13H32" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>'
            f'</svg>')

def svg_icon_check(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<circle cx="32" cy="32" r="24" fill="none" stroke="{color}" stroke-width="2.5"/>'
            f'<polyline points="20,32 28,42 44,22" fill="none" stroke="{color}" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/>'
            f'</svg>')

def svg_icon_layers(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<polygon points="32,8 56,22 32,36 8,22" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>'
            f'<polygon points="32,20 56,34 32,48 8,34" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round" opacity="0.7"/>'
            f'<polygon points="32,32 56,46 32,60 8,46" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round" opacity="0.4"/>'
            f'</svg>')

def svg_icon_refresh(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<path d="M48 16A22 22 0 0 0 12 26" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>'
            f'<path d="M16 48A22 22 0 0 0 52 38" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>'
            f'<polyline points="12,16 12,28 24,28" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>'
            f'<polyline points="52,48 52,36 40,36" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>'
            f'</svg>')

def svg_icon_rocket(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<path d="M32 4c-6 8-10 18-10 28 0 4 .5 7 1 10h18c.5-3 1-6 1-10 0-10-4-20-10-28z" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>'
            f'<circle cx="32" cy="26" r="4" fill="{color}" opacity="0.6"/>'
            f'<path d="M22 32c-6 2-10 6-10 10h10" fill="none" stroke="{color}" stroke-width="1.5"/>'
            f'<path d="M42 32c6 2 10 6 10 10H42" fill="none" stroke="{color}" stroke-width="1.5"/>'
            f'</svg>')

def svg_icon_target(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<circle cx="32" cy="32" r="26" fill="none" stroke="{color}" stroke-width="2"/>'
            f'<circle cx="32" cy="32" r="17" fill="none" stroke="{color}" stroke-width="2" opacity="0.7"/>'
            f'<circle cx="32" cy="32" r="8" fill="none" stroke="{color}" stroke-width="2" opacity="0.5"/>'
            f'<circle cx="32" cy="32" r="2.5" fill="{color}"/>'
            f'</svg>')

def svg_icon_database(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<ellipse cx="32" cy="14" rx="22" ry="8" fill="none" stroke="{color}" stroke-width="2"/>'
            f'<path d="M10 14v36c0 4.4 9.8 8 22 8s22-3.6 22-8V14" fill="none" stroke="{color}" stroke-width="2"/>'
            f'<path d="M10 26c0 4.4 9.8 8 22 8s22-3.6 22-8" fill="none" stroke="{color}" stroke-width="1.5" opacity="0.6"/>'
            f'<path d="M10 38c0 4.4 9.8 8 22 8s22-3.6 22-8" fill="none" stroke="{color}" stroke-width="1.5" opacity="0.6"/>'
            f'</svg>')

def svg_icon_network(size=64, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">'
            f'<circle cx="32" cy="14" r="7" fill="none" stroke="{color}" stroke-width="2.5"/>'
            f'<circle cx="14" cy="50" r="7" fill="none" stroke="{color}" stroke-width="2.5"/>'
            f'<circle cx="50" cy="50" r="7" fill="none" stroke="{color}" stroke-width="2.5"/>'
            f'<circle cx="32" cy="14" r="3" fill="{color}"/>'
            f'<circle cx="14" cy="50" r="3" fill="{color}"/>'
            f'<circle cx="50" cy="50" r="3" fill="{color}"/>'
            f'<line x1="32" y1="21" x2="14" y2="43" stroke="{color}" stroke-width="1.5"/>'
            f'<line x1="32" y1="21" x2="50" y2="43" stroke="{color}" stroke-width="1.5"/>'
            f'<line x1="21" y1="50" x2="43" y2="50" stroke="{color}" stroke-width="1.5"/>'
            f'</svg>')

def svg_chevron_flow(items, w=800, h=70, color="#7E00FF"):
    n = len(items)
    cw = w / n
    arrow = 14
    shapes = ""
    for i, lbl in enumerate(items):
        x = i * cw
        op = max(0.5, 1.0 - i * 0.12)
        shapes += (f'<polygon points="{x},0 {x+cw-arrow},0 {x+cw},{h/2} {x+cw-arrow},{h} {x},{h}"'
                   f' fill="{color}" opacity="{op}"/>')
        shapes += (f'<text x="{x+cw/2}" y="{h/2+5}" text-anchor="middle"'
                   f' font-family="Graphik, sans-serif" font-size="14"'
                   f' fill="white" font-weight="bold">{lbl}</text>')
    return f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}">{shapes}</svg>'

def _chevron_band_svg(labels, w=800, h=55):
    n = len(labels)
    sw = w / n
    td = h * 0.38
    alts = ["#5F0095", "#7E00FF"]
    polys = []
    for i, lbl in enumerate(labels):
        c = alts[i % 2]
        x0 = i * sw
        x1 = x0 + sw
        mid = h / 2
        pts = f"{x0},0 {x1-td},0 {x1},{mid} {x1-td},{h} {x0},{h}"
        polys.append(
            f'<polygon points="{pts}" fill="{c}"/>'
            f'<text x="{x0+sw/2}" y="{mid}" text-anchor="middle" dominant-baseline="middle"'
            f' fill="white" font-size="13" font-weight="bold"'
            f' font-family="Graphik, sans-serif">{lbl}</text>'
        )
    return f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {w} {h}">{"".join(polys)}</svg>'


# ── LAYOUT HELPERS ────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, border=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border is not None:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def add_tb(slide, x, y, w, h, text, size=14, bold=False, color=None,
           align=PP_ALIGN.LEFT, italic=False, wrap=True):
    if color is None:
        color = TEXT_BODY
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return tb

def set_slide_title(slide, title, breadcrumb):
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(28)
                p.font.color.rgb = BLACK
                p.font.name = FONT
        elif idx == 10:
            ph.text = ""
        elif idx == 11:
            ph.text = breadcrumb
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = PURPLE
                p.font.name = FONT

def set_lead(slide, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10:
            ph.text = text
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(18)
                p.font.color.rgb = BLACK
                p.font.name = FONT

layout_cover   = prs.slide_layouts[0]
layout_content = prs.slide_layouts[2]

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(layout_cover)
for ph in s1.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "Claude Code: Accelerating\nMeridian Claims Delivery"
        for p in ph.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = BLACK
            p.font.name = FONT
    elif idx == 1:
        ph.text = "A Delivery Operating Model — Evidenced, Not Proposed"
        for p in ph.text_frame.paragraphs:
            p.font.color.rgb = BLACK
            p.font.name = FONT
    elif idx == 12:
        ph.text = "RFP Technical Capability  |  AON Meridian Claims Copilot  |  August 2026"
        for p in ph.text_frame.paragraphs:
            p.font.color.rgb = BLACK
            p.font.name = FONT
    elif idx == 2:
        ph.text = "August 2026"
        for p in ph.text_frame.paragraphs:
            p.font.color.rgb = TEXT_SUB
            p.font.name = FONT
add_logo_to_cover(s1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1b: OUR PROPOSITION — Pattern E (Three-column hook)
# ═══════════════════════════════════════════════════════════════════
s_prop = prs.slides.add_slide(layout_content)
set_slide_title(s_prop, "We did not build a demo — we built the answer", "Our Proposition")
set_lead(s_prop, "Every bidder claims AI acceleration. We arrived with a running system, a skills library, and a delivery contract.")

prop_cols = [
    ("The Evidence",
     svg_icon_check,
     "Proof, not a proposal.",
     ["45 endpoints, 116 tests, BR-001 enforced — running today",
      "Built against your architecture — not illustrative of it",
      "FNOL resilience live: stop the backend, the claim still submits",
      "35/37 RFP features demonstrated in the proof-of-concept"],
     DEEP_PURPLE),
    ("The Operating\nModel",
     svg_icon_layers,
     "Standards encoded, not promised.",
     ["Skills encode your RFP requirements once, enforced on every commit",
      "25+ custom skills mapped to named NFRs and RFP exhibits",
      "CLAUDE.md carries context every session — zero architectural drift",
      "POC to production is a checklist, not a rewrite"],
     PURPLE),
    ("The Scale\nProof",
     svg_icon_network,
     "Enterprise-proven, not piloted.",
     ["Accenture × Anthropic: 200+ certified Skills in production",
      "100,000+ users deployed — proven enterprise methodology",
      "Skills COE: Build → Validate → Govern → Distribute lifecycle",
      "This POC runs on the framework already proven at scale"],
     DEEP_PURPLE),
]

N_pc = len(prop_cols)
col_w_pc = (CW - (N_pc - 1) * 0.15) / N_pc

for i, (title, icon_fn, subtitle, bullets, hdr_color) in enumerate(prop_cols):
    cx = ML + i * (col_w_pc + 0.15)
    add_rect(s_prop, cx, CY, col_w_pc, AH, fill=WHITE, border=LIGHT_GRAY)
    add_rect(s_prop, cx, CY, col_w_pc, 0.05, fill=hdr_color)
    add_svg(s_prop, icon_fn(36, "#7E00FF"), cx + (col_w_pc - 0.42) / 2, CY + 0.15, 0.42, 0.42, 256)
    add_tb(s_prop, cx + 0.10, CY + 0.65, col_w_pc - 0.20, 0.60,
           title, size=14, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_tb(s_prop, cx + 0.10, CY + 1.30, col_w_pc - 0.20, 0.30,
           subtitle, size=14, italic=True, color=BLACK)
    tb_bp = s_prop.shapes.add_textbox(Inches(cx + 0.12), Inches(CY + 1.65),
                                      Inches(col_w_pc - 0.24), Inches(AH - 1.75))
    tf_bp = tb_bp.text_frame
    tf_bp.word_wrap = True
    for j, line in enumerate(bullets):
        p = tf_bp.add_paragraph() if j > 0 else tf_bp.paragraphs[0]
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = f"• {line}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT

add_gt_to_slide(s_prop)
set_footer(s_prop)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(layout_content)
for ph in s2.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "Agenda"
        for p in ph.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(28)
            p.font.color.rgb = BLACK
            p.font.name = FONT
    elif idx == 10:
        ph.text = ""
    elif idx == 11:
        ph.text = "Overview"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = PURPLE
            p.font.name = FONT

agenda_items = [
    "The Proof Is Already Running",
    "How Claude Accelerates Delivery",
    "Platform Skills: Day-One Gates",
    "Custom Compliance Skills",
    "FastAPI and Python Optimisers",
    "Claims Domain Skills",
    "RFP-Specific Skills: Closing Gaps",
    "Skills in the Delivery Cycle",
    "Agents: Parallel Execution Engine",
    "CLAUDE.md: The Machine-Readable Contract",
    "Delivery Acceleration Model",
    "What We Bring Forward",
]
tb_ag = s2.shapes.add_textbox(Inches(ML), Inches(CY), Inches(CW), Inches(AH))
tf_ag = tb_ag.text_frame
tf_ag.word_wrap = True
for i, item in enumerate(agenda_items):
    p = tf_ag.add_paragraph() if i > 0 else tf_ag.paragraphs[0]
    p.space_before = Pt(6)
    r = p.add_run()
    r.text = f"{i + 1}.   {item}"
    r.font.size = Pt(18)
    r.font.color.rgb = TEXT_BODY
    r.font.name = FONT
add_gt_to_slide(s2)
set_footer(s2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: THE PROOF  — Pattern F (Numbered summary rows)
# ═══════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(layout_content)
set_slide_title(s3, "This demo was built — not scaffolded — by Claude", "The Proof")
set_lead(s3, "45 endpoints · 12 routers · 116 tests · Webpack 5 MFE shell. Built end-to-end with Claude Code as the delivery partner.")

proof_rows = [
    ("01", "Backend API",       "45 endpoints  ·  12 routers  ·  116 passing tests  ·  Python FastAPI + Pydantic v2",  svg_icon_network),
    ("02", "Frontend",          "Webpack 5 Module Federation  ·  7 role-scoped screens  ·  7 demo personas  ·  5 locales",  svg_icon_code),
    ("03", "Security",          "BR-001 org scope on every route  ·  3-gate document model  ·  JWT HS256  →  RS256 ready",  svg_icon_shield),
    ("04", "Production Paths",  "SQLite → MySQL  ·  Mock auth → Okta PKCE  ·  uui-stub → @aon/united-ui  ·  All documented",  svg_icon_layers),
]
msg_h = 0.60
n_rows = len(proof_rows)
gap = 0.10
item_h = (AH - msg_h - (n_rows + 1) * gap) / n_rows
badge_w = 0.52

y = CY + gap
for i, (num, label, detail, icon_fn) in enumerate(proof_rows):
    add_rect(s3, ML, y, badge_w, item_h, fill=DEEP_PURPLE)
    add_tb(s3, ML, y + (item_h - 0.30) / 2, badge_w, 0.30,
           num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card_x = ML + badge_w + 0.08
    card_w = CW - badge_w - 0.08
    add_rect(s3, card_x, y, card_w, item_h, fill=BG_LIGHT)
    add_tb(s3, card_x + 0.15, y + 0.10, 2.20, item_h - 0.15,
           label, size=14, bold=True, color=DEEP_PURPLE)
    add_tb(s3, card_x + 2.45, y + 0.10, card_w - 2.75, item_h - 0.15,
           detail, size=14, color=TEXT_BODY)
    icon_x = card_x + card_w - 0.55
    icon_y = y + (item_h - 0.38) / 2
    add_svg(s3, icon_fn(32, "#7E00FF"), icon_x, icon_y, 0.38, 0.38, 256)
    y += item_h + gap

msg_y = CY + gap + n_rows * (item_h + gap)
add_rect(s3, ML, msg_y, CW, msg_h, fill=LP_BG, border=PURPLE)
add_tb(s3, ML + 0.20, msg_y + 0.12, CW - 0.40, msg_h - 0.15,
       "▸  Not a prototype — a production-pattern system built at prototype speed.",
       size=14, bold=True, color=DEEP_PURPLE)
add_gt_to_slide(s3)
set_footer(s3)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3b: HOW CLAUDE ACCELERATES DELIVERY — 3-row before/after compression
# ═══════════════════════════════════════════════════════════════════
s_acc = prs.slides.add_slide(layout_content)
set_slide_title(s_acc, "Three loops Claude permanently compresses in delivery", "How Claude Accelerates Delivery")
set_lead(s_acc, "Speed is not the claim — loop compression is. These three loops determine the total delivery cost of any engagement.")

loop_data = [
    ("Quality\nGate Loop",
     DEEP_PURPLE,
     svg_icon_shield,
     "Security and compliance reviews run post-merge — caught by a separate team, too late to fix cheaply",
     "30-second diff-targeted scan on every commit — /security-review, /sonar-gate, /rbac-matrix run before merge",
     "Defects caught at zero rework cost"),
    ("Context\nRebuild Loop",
     PURPLE,
     svg_icon_document,
     "Each engineer re-derives architecture from docs and tribal knowledge — context and standards drift across sprints",
     "CLAUDE.md loads every session — 51 NFRs, security rules, swap paths, folder nav inherited automatically",
     "Senior-architect context from commit one"),
    ("Research-\nBuild Loop",
     DEEP_PURPLE,
     svg_icon_search,
     "Cross-cutting changes block on understanding — investigate first, then implement; research and build run in sequence",
     "Explore + Plan agents run concurrently — research proceeds in background while implementation starts",
     "40–60% compression on cross-cutting changes"),
]

label_w_l  = 1.80
gap_cells  = 0.03
avail_w_l  = CW - label_w_l - 2 * gap_cells
before_w_l = avail_w_l * 0.43
after_w_l  = avail_w_l - before_w_l
n_lp       = len(loop_data)
gap_r_l    = 0.10
row_h_l    = (AH - (n_lp - 1) * gap_r_l) / n_lp

y = CY
for i, (loop_name, lcolor, icon_fn, before, after, impact) in enumerate(loop_data):
    add_rect(s_acc, ML, y, label_w_l, row_h_l, fill=lcolor)
    add_svg(s_acc, icon_fn(28, "#FFFFFF"), ML + (label_w_l - 0.30) / 2, y + 0.12, 0.30, 0.30, 256)
    add_tb(s_acc, ML + 0.10, y + 0.48, label_w_l - 0.20, 0.70,
           loop_name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bx = ML + label_w_l + gap_cells
    add_rect(s_acc, bx, y, before_w_l, row_h_l, fill=BG_LIGHT)
    add_tb(s_acc, bx + 0.10, y + 0.10, before_w_l - 0.20, 0.20,
           "BEFORE", size=14, bold=True, color=TEXT_SUB, italic=True)
    add_tb(s_acc, bx + 0.10, y + 0.34, before_w_l - 0.20, row_h_l - 0.44,
           before, size=14, color=TEXT_BODY)
    ax = bx + before_w_l + gap_cells
    add_rect(s_acc, ax, y, after_w_l, row_h_l, fill=LP_BG)
    add_tb(s_acc, ax + 0.10, y + 0.10, after_w_l - 0.20, 0.20,
           "AFTER", size=14, bold=True, color=PURPLE)
    add_tb(s_acc, ax + 0.10, y + 0.34, after_w_l - 0.20, row_h_l - 0.80,
           after, size=14, color=DEEP_PURPLE)
    stat_y = y + row_h_l - 0.40
    add_rect(s_acc, ax, stat_y, after_w_l, 0.40, fill=lcolor)
    add_tb(s_acc, ax + 0.12, stat_y + 0.08, after_w_l - 0.24, 0.26,
           f"► {impact}", size=14, bold=True, color=WHITE)
    y += row_h_l + (gap_r_l if i < n_lp - 1 else 0)

add_gt_to_slide(s_acc)
set_footer(s_acc)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: PLATFORM SKILLS  — Pattern F (Numbered row-cards, same as Proof)
# ═══════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(layout_content)
set_slide_title(s4, "Five platform skills gate every change from day one", "Platform Skills")
set_lead(s4, "These skills ship with Claude Code and are active from project setup — no authoring required.")

platform_rows = [
    ("/code-review",     "Correctness on the diff",        "Scans changed lines for bugs and simplification — not the whole file; engineers see exactly what they introduced",    svg_icon_code),
    ("/security-review", "OWASP top-10 on every commit",   "Runs OWASP analysis against changed lines; catches injection, auth bypass, and secrets before they are merged",       svg_icon_shield),
    ("/simplify",        "Complexity gate before merge",   "Strips premature abstractions and generalisations; pushes back before over-engineered code reaches the codebase",     svg_icon_check),
    ("/run",             "Golden path before done",        "Launches the full stack and exercises the browser path — task is not reported complete until this passes",             svg_icon_rocket),
    ("/init",            "Architecture on day one",        "Generates CLAUDE.md from the live codebase — security rules, folder nav, and production swap paths encoded once",     svg_icon_document),
]

n_rows_p = len(platform_rows)
gap_p    = 0.10
row_h_p  = (AH - (n_rows_p - 1) * gap_p) / n_rows_p
badge_w_p = 0.52

y = CY
for i, (skill, label, desc, icon_fn) in enumerate(platform_rows):
    # Numbered badge
    add_rect(s4, ML, y, badge_w_p, row_h_p, fill=DEEP_PURPLE)
    add_tb(s4, ML, y + (row_h_p - 0.30) / 2, badge_w_p, 0.30,
           f"{i+1:02d}", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Content card
    card_x4 = ML + badge_w_p + 0.08
    card_w4  = CW - badge_w_p - 0.08
    add_rect(s4, card_x4, y, card_w4, row_h_p, fill=BG_LIGHT)
    # Skill name column
    add_tb(s4, card_x4 + 0.15, y + 0.10, 2.10, row_h_p - 0.14,
           skill, size=14, bold=True, color=DEEP_PURPLE)
    # Purple subtitle
    add_tb(s4, card_x4 + 2.38, y + 0.10, card_w4 - 3.10, 0.26,
           label, size=14, bold=True, color=PURPLE)
    # Description
    add_tb(s4, card_x4 + 2.38, y + 0.38, card_w4 - 3.10, row_h_p - 0.48,
           desc, size=14, color=TEXT_BODY)
    # Icon
    add_svg(s4, icon_fn(32, "#7E00FF"),
            card_x4 + card_w4 - 0.55, y + (row_h_p - 0.38) / 2, 0.38, 0.38, 256)
    y += row_h_p + (gap_p if i < n_rows_p - 1 else 0)

add_gt_to_slide(s4)
set_footer(s4)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: COMPLIANCE SKILLS  — Pattern M (Value Proposition Rows)
# ═══════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(layout_content)
set_slide_title(s5, "Compliance skills enforce your enterprise standards", "Custom Skills: Compliance")
set_lead(s5, "Three custom skills run SonarQube, Fortify, and OWASP scans against the diff — not the whole codebase — on every change.")

compliance_rows = [
    ("SonarQube\nGate",
     "Diff-targeted quality analysis — not a full project scan",
     ["Maps SonarQube findings to the changed lines only — engineers see exactly what they introduced",
      "Understands your configured rule profile; blocks the session if your quality gate fails",
      "Works against your existing SonarQube server — zero new infrastructure"]),
    ("Fortify\n+ OWASP",
     "SAST scan and dependency CVE check in one targeted pass",
     ["Fortify FPR output parsed to critical and high findings with exact file and line citations",
      "OWASP Dependency-Check covers requirements.txt and package.json including transitive CVEs",
      "Both resolve in under 90 seconds on the Claims API stack"]),
    ("License\nCompliance",
     "GPL contamination caught at install time — not at legal review",
     ["Validates all new dependencies against an approved SPDX list before they are committed",
      "Covers transitive dependencies — not just the direct install",
      "Zero new tooling; runs against npm and pip install output already in the terminal"]),
]

LABEL_W = 1.80
gap = 0.10
row_h = (AH - (len(compliance_rows) - 1) * gap) / len(compliance_rows)
icons_comp = [svg_icon_search, svg_icon_shield, svg_icon_document]

y = CY
for i, (label, lead, bullets) in enumerate(compliance_rows):
    # Purple label block
    add_rect(s5, ML, y, LABEL_W, row_h, fill=DEEP_PURPLE if i % 2 == 0 else PURPLE)
    add_tb(s5, ML + 0.05, y + row_h * 0.30, LABEL_W - 0.10, row_h * 0.40,
           label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Content area
    card_x = ML + LABEL_W + 0.08
    card_w = CW - LABEL_W - 0.08
    add_rect(s5, card_x, y, card_w, row_h, fill=BG_LIGHT)
    # Icon
    add_svg(s5, icons_comp[i](28, "#7E00FF"), card_x + 0.12, y + 0.12, 0.32, 0.32, 256)
    # Lead text
    add_tb(s5, card_x + 0.55, y + 0.10, card_w - 0.70, 0.30,
           lead, size=14, bold=True, color=BLACK)
    # Bullets
    tb_b = s5.shapes.add_textbox(Inches(card_x + 0.20), Inches(y + 0.45),
                                  Inches(card_w - 0.35), Inches(row_h - 0.55))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    for j, bullet in enumerate(bullets):
        p = tf_b.add_paragraph() if j > 0 else tf_b.paragraphs[0]
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = f"- {bullet}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT
    y += row_h + (gap if i < len(compliance_rows) - 1 else 0)

add_gt_to_slide(s5)
set_footer(s5)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: FASTAPI & PYTHON  — Pattern V (Column Chevron + Icon Cards)
# ═══════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(layout_content)
set_slide_title(s6, "FastAPI optimisers stop backend drift before merge", "Custom Skills: Optimisers")
set_lead(s6, "Four custom skills targeting the specific failure modes of a FastAPI / Pydantic v2 stack under production load.")

v_labels  = ["FastAPI Audit", "N+1 Guard", "Pydantic v2", "Async Audit"]
v_icons   = [svg_icon_gear, svg_icon_database, svg_icon_refresh, svg_icon_rocket]
v_bullets = [
    ["Checks every modified route for missing current_scope dependency",
     "Flags org_node read from params/headers (BR-001 breach)",
     "Catches routes that need a server restart to register"],
    ["Detects query-per-row loops before they reach production",
     "60 claims in dev is invisible; 600,000 claims pages on-call at 2 a.m.",
     "Covers SQLite dev paths and MySQL production paths equally"],
    ["Enforces Pydantic v2 idioms — no .dict(), no orm_mode",
     "Any v1 pattern silently degrades performance",
     "Will break on the next minor release if not caught now"],
    ["Flags synchronous blocking calls inside async def route handlers",
     "One blocking call serialises the entire event loop under load",
     "Catches requests library, file I/O, and time.sleep misuse"],
]

N = len(v_labels)
chev_h = 0.55
gap_cv = 0.10
card_y = CY + chev_h + gap_cv
card_h = AH - chev_h - gap_cv
card_w = (CW - (N - 1) * gap_cv) / N

add_svg(s6, _chevron_band_svg(v_labels, w=800, h=55),
        ML, CY, CW, chev_h, pw=800)

for i, (lbl, icon_fn, bullets) in enumerate(zip(v_labels, v_icons, v_bullets)):
    cx = ML + i * (card_w + gap_cv)
    add_rect(s6, cx, card_y, card_w, card_h, fill=WHITE, border=LIGHT_GRAY)
    # Purple top bar
    add_rect(s6, cx, card_y, card_w, 0.05, fill=PURPLE)
    # Icon top-center
    icon_sz = 0.35
    add_svg(s6, icon_fn(32, "#7E00FF"),
            cx + (card_w - icon_sz) / 2, card_y + 0.12, icon_sz, icon_sz, 256)
    # Title
    add_tb(s6, cx + 0.08, card_y + 0.55, card_w - 0.16, 0.36,
           lbl, size=14, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    # Bullets
    tb_v = s6.shapes.add_textbox(Inches(cx + 0.12), Inches(card_y + 0.98),
                                   Inches(card_w - 0.24), Inches(card_h - 1.05))
    tf_v = tb_v.text_frame
    tf_v.word_wrap = True
    for j, line in enumerate(bullets):
        p = tf_v.add_paragraph() if j > 0 else tf_v.paragraphs[0]
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = f"• {line}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT

add_gt_to_slide(s6)
set_footer(s6)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: DOMAIN SKILLS  — Pattern B (2×2 Grid)
# ═══════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(layout_content)
set_slide_title(s7, "Domain skills keep the demo and production coherent", "Custom Skills: Domain")
set_lead(s7, "Four skills targeting the unique risks of a claims demo demoed to SMEs and an eventual production swap.")

domain_cells = [
    ("/scope-audit",      svg_icon_shield,
     ["Scans every modified route for BR-001 compliance",
      "Verifies current_scope applied, org_node from JWT only",
      "Checks 403 (not 404) returned and audit log written on deny"]),
    ("/seed-coherence",   svg_icon_target,
     ["Validates scenario binding: product → cause → consequence → narrative",
      "Catches 'Ransomware on Motor Fleet / Escape of Water' before the demo",
      "References LOSS_SCENARIOS taxonomy in seed.py automatically"]),
    ("/prod-readiness",   svg_icon_layers,
     ["Scans for all POC stubs: uui-stub, mock auth, SQLite, HS256",
      "Reports open production swap paths against CLAUDE.md targets",
      "Turns the readiness audit from a manual checklist into a 10-second command"]),
    ("/migration-safe",   svg_icon_database,
     ["Validates that every schema change is additive-only",
      "No DROP, no ALTER removing columns, no NOT NULL without default",
      "Prevents a production MySQL InnoDB incident from a dev shortcut"]),
]

cell_w = (CW - 0.20) / 2
cell_h = (AH - 0.20) / 2
gap_b = 0.20

for idx, (title, icon_fn, bullets) in enumerate(domain_cells):
    col = idx % 2
    row = idx // 2
    cx = ML + col * (cell_w + gap_b)
    cy = CY + row * (cell_h + 0.10)
    add_rect(s7, cx, cy, cell_w, cell_h, fill=WHITE, border=LIGHT_GRAY)
    add_rect(s7, cx, cy, cell_w, 0.05, fill=PURPLE)
    # Icon
    add_svg(s7, icon_fn(28, "#7E00FF"), cx + 0.15, cy + 0.12, 0.34, 0.34, 256)
    # Title
    add_tb(s7, cx + 0.58, cy + 0.10, cell_w - 0.70, 0.38,
           title, size=14, bold=True, color=DEEP_PURPLE)
    # Bullets
    tb_d = s7.shapes.add_textbox(Inches(cx + 0.15), Inches(cy + 0.55),
                                   Inches(cell_w - 0.30), Inches(cell_h - 0.65))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    for j, line in enumerate(bullets):
        p = tf_d.add_paragraph() if j > 0 else tf_d.paragraphs[0]
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = f"• {line}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT

add_gt_to_slide(s7)
set_footer(s7)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: RFP-SPECIFIC SKILLS  — Pattern C (3-col Table)
# ═══════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(layout_content)
set_slide_title(s8, "New skills target every material gap in the AON RFP", "RFP-Specific Skills")
set_lead(s8, "Each skill maps to a named NFR, open item, or RFP exhibit — not a best practice, a specific clause.")

rfp_skills = [
    ("/wcag-audit",        "NFR-48/49/50", "Per-release WCAG 2.2 AA record via axe-core across all 7 screens — not a one-time audit"),
    ("/datadog-trace",     "NFR-51 Must",  "Every notification write verified to have a Datadog event — replaces the TODO comment"),
    ("/fnol-outbox-check", "NFR-37 Must",  "Outbox-first FNOL write enforced; prevents any sprint bypassing the resilience pattern"),
    ("/okta-contract",     "OI-26 / R-14", "JWT claims shape, JWKS path, and org_node source validated — never reads from params"),
    ("/perf-budget",       "NFR-26/27",    "45 endpoints classified by complexity band; blocks any exceeding their response cap"),
    ("/rbac-matrix",       "Exhibit 5",    "7-persona privilege matrix + multi-row AND logic verified every sprint"),
    ("/nfr-matrix",        "All 51 NFRs",  "Auto-generates the submission compliance matrix from the codebase each sprint"),
    ("/snyk-gate",         "Aon DevOps",   "Container and Terraform IaC CVE scan — complements /sonar-gate + /owasp-dependency"),
]

n_cols_r = 2
n_rows_r = 4
cell_gap_r = 0.10
cell_w_r = (CW - (n_cols_r - 1) * cell_gap_r) / n_cols_r
cell_h_r = (AH - (n_rows_r - 1) * cell_gap_r) / n_rows_r
badge_w_r = 1.50

for idx, (skill, rfp_ref, desc) in enumerate(rfp_skills):
    row = idx // n_cols_r
    col = idx % n_cols_r
    cx = ML + col * (cell_w_r + cell_gap_r)
    cy = CY + row * (cell_h_r + cell_gap_r)
    # Card
    add_rect(s8, cx, cy, cell_w_r, cell_h_r, fill=BG_LIGHT)
    add_rect(s8, cx, cy, cell_w_r, 0.05, fill=PURPLE)  # top accent
    # Skill name (left)
    add_tb(s8, cx + 0.15, cy + 0.12, 2.20, 0.30,
           skill, size=14, bold=True, color=DEEP_PURPLE)
    # NFR badge (right, PURPLE rect + white text)
    badge_x = cx + cell_w_r - badge_w_r - 0.12
    add_rect(s8, badge_x, cy + 0.12, badge_w_r, 0.30, fill=PURPLE)
    add_tb(s8, badge_x, cy + 0.12, badge_w_r, 0.30,
           rfp_ref, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Description
    add_tb(s8, cx + 0.15, cy + 0.48, cell_w_r - 0.30, cell_h_r - 0.58,
           desc, size=14, color=TEXT_BODY)

add_gt_to_slide(s8)
set_footer(s8)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: SKILLS IN THE DELIVERY CYCLE  — Custom Swimlane
# ═══════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(layout_content)
set_slide_title(s9, "Every phase of delivery has a skill that enforces it", "Skills in the Delivery Cycle")
set_lead(s9, "Skills embed in Sprint 0 setup, every commit, every PR gate, and every release — not a review layer bolted on at the end.")

sdlc_phases = ["Sprint 0", "Every Change", "PR Gate", "Sprint Close", "Release"]
sdlc_skills = [
    # Sprint 0
    [("/init",              "P"), ("/okta-contract",     "C"),
     ("/seed-coherence",    "C"), ("/country-config",    "C")],
    # Every Change
    [("/code-review",       "P"), ("/security-review",   "P"),
     ("/scope-audit",       "C"), ("/fastapi-audit",     "C"),
     ("/rbac-matrix",       "C"), ("/fnol-outbox-check", "C")],
    # PR Gate
    [("/sonar-gate",        "C"), ("/snyk-gate",         "C"),
     ("/fortify-scan",      "C"), ("/owasp-dependency",  "C"),
     ("/audit-completeness","C"), ("/i18n-parity",       "C")],
    # Sprint Close
    [("/wcag-audit",        "C"), ("/perf-budget",       "C"),
     ("/migration-safe",    "C"), ("/nfr-matrix",        "C"),
     ("/uui-compliance",    "C")],
    # Release
    [("/datadog-trace",     "C"), ("/prod-readiness",    "C"),
     ("/mfe-deploy-check",  "C"), ("/acord-schema",      "C")],
]

chev_h9  = 0.55
gap9     = 0.08
card_top = CY + chev_h9 + gap9
legend_h = 0.30
card_h9  = BY - legend_h - gap9 - card_top
N9       = len(sdlc_phases)
col_gap9 = 0.08
col_w9   = (CW - (N9 - 1) * col_gap9) / N9

# Phase chevron band
add_svg(s9, _chevron_band_svg(sdlc_phases, w=800, h=55),
        ML, CY, CW, chev_h9, pw=800)

# Column cards + skill names
body_h_pts = (card_h9 - 0.34) * 72  # subtract card header
for ci, (phase, skills) in enumerate(zip(sdlc_phases, sdlc_skills)):
    cx = ML + ci * (col_w9 + col_gap9)
    # Card background (no repeated header — chevron band above already labels phases)
    hdr_color = DEEP_PURPLE if ci % 2 == 0 else PURPLE
    add_rect(s9, cx, card_top, col_w9, card_h9, fill=WHITE, border=None)
    add_rect(s9, cx, card_top, col_w9, 0.06, fill=hdr_color)  # thin top accent only
    # Dynamic spacing to fill ~78% of card
    n = len(skills)
    text_pts = n * 16
    gap_pts  = 0.78 * body_h_pts - text_pts
    top_pad  = 10
    sb_pts   = max(6, min(48, (gap_pts - top_pad) / max(n - 1, 1)))
    # Skill names
    tb9 = s9.shapes.add_textbox(
        Inches(cx + 0.10), Inches(card_top + 0.14),
        Inches(col_w9 - 0.20), Inches(card_h9 - 0.18)
    )
    tf9 = tb9.text_frame
    tf9.word_wrap = True
    for si, (skill_name, category) in enumerate(skills):
        p = tf9.add_paragraph() if si > 0 else tf9.paragraphs[0]
        p.space_before = Pt(int(sb_pts) if si > 0 else 0)
        r = p.add_run()
        r.text = skill_name
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = DEEP_PURPLE if category == "P" else PURPLE
        r.font.name = FONT

# Legend
leg_y = BY - legend_h
add_rect(s9, ML, leg_y, CW, legend_h, fill=LP_BG)
add_tb(s9, ML + 0.20, leg_y + 0.05, 5.50, 0.24,
       "■  Platform — ships with Claude Code", size=14, bold=False, color=DEEP_PURPLE)
add_tb(s9, ML + 6.00, leg_y + 0.05, 5.50, 0.24,
       "■  Custom — built for AON Meridian", size=14, bold=False, color=PURPLE)

add_svg(s9, svg_icon_layers(28, "#7E00FF"), ML + CW - 0.50, leg_y - 0.48, 0.36, 0.36, 256)
add_gt_to_slide(s9)
set_footer(s9)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: AGENTS  — Pattern E (Three-Column)
# ═══════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(layout_content)
set_slide_title(s10, "Parallel agents compress research and build cycles", "Agents")
set_lead(s10, "Four specialist sub-agents run concurrently — research, architecture, review, and execution in parallel, not sequence.")

agent_cols = [
    ("Explore",  svg_icon_search,
     "Fast read-only codebase search.",
     ["Finds symbols and traces call chains across 700+ files",
      "Answers 'where is BR-001 enforced' without touching anything",
      "Runs in background — primary session keeps building"]),
    ("Plan",     svg_icon_brain,
     "Architecture before a line is written.",
     ["Identifies blast radius for every proposed change",
      "Surfaces ADR candidates and boundary violations",
      "Returns a step-by-step plan with verify checkpoints"]),
    ("Review\n+ Execute", svg_icon_network,
     "Independent review and multi-step orchestration.",
     ["code-reviewer gives second opinion without seeing author reasoning",
      "general-purpose runs grep + read + analyse in one background pass",
      "Parallel execution cuts the research-build loop by 40–60%"]),
]

N3 = len(agent_cols)
col_w = (CW - (N3 - 1) * 0.15) / N3
col_h = AH
gap_e = 0.15

for i, (title, icon_fn, subtitle, bullets) in enumerate(agent_cols):
    cx = ML + i * (col_w + gap_e)
    add_rect(s10, cx, CY, col_w, col_h, fill=WHITE, border=LIGHT_GRAY)
    add_rect(s10, cx, CY, col_w, 0.05, fill=DEEP_PURPLE if i == 0 else PURPLE)
    # Icon
    add_svg(s10, icon_fn(36, "#7E00FF"), cx + (col_w - 0.42) / 2, CY + 0.15, 0.42, 0.42, 256)
    # Title
    add_tb(s10, cx + 0.10, CY + 0.65, col_w - 0.20, 0.55,
           title, size=14, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    # Subtitle
    add_tb(s10, cx + 0.10, CY + 1.22, col_w - 0.20, 0.38,
           subtitle, size=14, bold=True, color=BLACK, italic=True)
    # Bullets
    tb_e = s10.shapes.add_textbox(Inches(cx + 0.12), Inches(CY + 1.65),
                                   Inches(col_w - 0.24), Inches(col_h - 1.75))
    tf_e = tb_e.text_frame
    tf_e.word_wrap = True
    for j, line in enumerate(bullets):
        p = tf_e.add_paragraph() if j > 0 else tf_e.paragraphs[0]
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = f"• {line}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT

add_gt_to_slide(s10)
set_footer(s10)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: CLAUDE.MD  — Pattern J (Proposal Overview)
# ═══════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(layout_content)
set_slide_title(s11, "CLAUDE.md encodes the contract every session loads", "CLAUDE.md")
set_lead(s11, "A machine-readable project file that ensures every AI-assisted change carries the same context a senior architect would bring.")

j_rows = [
    ("Security\nRules",
     "Four non-negotiable security constraints enforced automatically",
     [("BR-001", "org_node from JWT only — never from query params, headers, or body on any data route"),
      ("Pillar 1", "Three sequential document gates: audience → security attribute → provenance — all in sdms_proxy.py"),
      ("ecm_reference", "Must never appear in any API response — checked on every document-related endpoint change"),
      ("Auth modes", "Mock HS256 is development-only; AUTH_MODE=okta delegates to Okta RS256 JWKS in production")]),
    ("Architecture\nObligations",
     "Update rules that keep Architecture.md current — not aspirational",
     [("Doc sync", "Architecture.md updated in the same change that adds a table, router, endpoint, or screen"),
      ("Swap paths", "Five production targets documented: MySQL, Okta, @aon/united-ui, Reltio, ECM FileNet"),
      ("Folder index", "Navigate 700 files via the CLAUDE.md directory table — not by scanning node_modules"),
      ("Migrations", "Additive-only enforced via ADDITIVE_COLUMNS in schema.py — no silent ALTER or DROP")]),
]

highlight_h = 0.65
j_gap = 0.10
row_h_j = (AH - highlight_h - 3 * j_gap) / 2
JLABEL_W = 1.80

y_j = CY
for i, (label, lead, items) in enumerate(j_rows):
    add_rect(s11, ML, y_j, JLABEL_W, row_h_j, fill=DEEP_PURPLE if i == 0 else PURPLE)
    add_tb(s11, ML + 0.05, y_j + row_h_j * 0.30, JLABEL_W - 0.10, row_h_j * 0.40,
           label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Content card
    card_x = ML + JLABEL_W + 0.08
    card_w = CW - JLABEL_W - 0.08
    add_rect(s11, card_x, y_j, card_w, row_h_j, fill=BG_LIGHT)
    # Lead
    add_tb(s11, card_x + 0.15, y_j + 0.10, card_w - 0.30, 0.30,
           f"• {lead}", size=14, bold=True, color=BLACK)
    # Items
    tb_j = s11.shapes.add_textbox(Inches(card_x + 0.15), Inches(y_j + 0.45),
                                   Inches(card_w - 0.30), Inches(row_h_j - 0.55))
    tf_j = tb_j.text_frame
    tf_j.word_wrap = True
    for k, (tag, desc) in enumerate(items):
        p = tf_j.add_paragraph() if k > 0 else tf_j.paragraphs[0]
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = f"{tag}: "
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = PURPLE
        r1.font.name = FONT
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = TEXT_BODY
        r2.font.name = FONT
    y_j += row_h_j + j_gap

# Highlight bar
hl_y = y_j + j_gap * 0.5
add_rect(s11, ML, hl_y, CW, highlight_h, fill=LP_BG, border=PURPLE)
add_tb(s11, ML + 0.20, hl_y + 0.12, CW - 0.40, highlight_h - 0.15,
       "▸  The CLAUDE.md is not documentation for humans — it is a machine-readable contract every session inherits.",
       size=14, bold=True, color=DEEP_PURPLE)
add_svg(s11, svg_icon_document(28, "#7E00FF"), ML + CW - 0.52, CY + 0.08, 0.38, 0.38, 256)

add_gt_to_slide(s11)
set_footer(s11)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: DELIVERY ACCELERATION  — Pattern I (Chevron Process)
# ═══════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(layout_content)
set_slide_title(s12, "Four delivery loops compressed to minutes not days", "Delivery Model")
set_lead(s12, "On the Claims Copilot POC, Claude Code compressed four loops that consume 60–70% of engineering time on a typical project.")

chevron_labels = ["Encode", "Scan", "Execute", "Verify"]
chev_svg_h = 0.65
card_gap = 0.12
card_y_i = CY + chev_svg_h + card_gap
card_h_i = AH - chev_svg_h - card_gap
N4 = len(chevron_labels)
card_w_i = (CW - (N4 - 1) * card_gap) / N4

add_svg(s12, svg_chevron_flow(chevron_labels, w=800, h=65, color="#7E00FF"),
        ML, CY, CW, chev_svg_h, pw=800)

card_details = [
    ("Standards written once",
     ["Security rules, production targets, folder nav — all in CLAUDE.md",
      "Every engineer and session inherits the same context",
      "Written once in Sprint 0; applies for the full engagement"]),
    ("Diff-targeted quality gates",
     ["/security-review catches scope leaks in 30 seconds",
      "/sonar-gate and /fortify-scan run against changed lines only",
      "No full-project scan; no false-positive noise from legacy code"]),
    ("Parallel implementation",
     ["Cross-file changes tracked as a checklist — not left to memory",
      "Agents run concurrently; research and build do not block each other",
      "Architecture.md updated in the same pass as the code change"]),
    ("Golden path verified",
     ["/run exercises the app before the task is reported done",
      "116 tests pass before any change is considered complete",
      "Production swap paths re-checked against CLAUDE.md targets"]),
]

icons_i = [svg_icon_document, svg_icon_shield, svg_icon_network, svg_icon_check]

for i, (subtitle, bullets) in enumerate(card_details):
    cx = ML + i * (card_w_i + card_gap)
    add_rect(s12, cx, card_y_i, card_w_i, card_h_i, fill=WHITE, border=LIGHT_GRAY)
    # Icon
    add_svg(s12, icons_i[i](28, "#7E00FF"),
            cx + (card_w_i - 0.34) / 2, card_y_i + 0.12, 0.34, 0.34, 256)
    # Subtitle
    add_tb(s12, cx + 0.10, card_y_i + 0.55, card_w_i - 0.20, 0.40,
           subtitle, size=14, bold=True, color=DEEP_PURPLE, align=PP_ALIGN.CENTER)
    # Bullets
    tb_i = s12.shapes.add_textbox(Inches(cx + 0.12), Inches(card_y_i + 1.00),
                                    Inches(card_w_i - 0.24), Inches(card_h_i - 1.10))
    tf_i = tb_i.text_frame
    tf_i.word_wrap = True
    for j, line in enumerate(bullets):
        p = tf_i.add_paragraph() if j > 0 else tf_i.paragraphs[0]
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = f"• {line}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT

add_gt_to_slide(s12)
set_footer(s12)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12b: EFFICIENCY DRIVERS — 5 horizontal rows
# ═══════════════════════════════════════════════════════════════════
s_eff = prs.slides.add_slide(layout_content)
set_slide_title(s_eff, "Five drivers that deliver 40–55% AD efficiency gain", "Efficiency and Value Realization")
set_lead(s_eff, "Each driver is live in the POC — gains compound from Sprint 0 and track Lincoln Financial's own 40–55% AD productivity target from the AI Readiness RFI.")

n_eff    = 5
gap_eff  = 0.07
row_h_e  = (AH - (n_eff - 1) * gap_eff) / n_eff   # 0.984"
drv_w    = 2.40
mech_w   = 6.80
imp_w    = CW - drv_w - mech_w                     # 3.30"
mech_x_e = ML + drv_w
imp_x_e  = mech_x_e + mech_w

eff_drivers = [
    (svg_icon_code,    DEEP_PURPLE, "AI Assisted\nDevelopment",
     "25+ custom skills on every commit — /code-review, /security-review, /scope-audit, /rbac-matrix",
     "CLAUDE.md loaded every session — 51 NFRs, security rules, and swap paths inherited automatically",
     "30–35%", "developer productivity lift", "LFG AI Readiness Q1"),
    (svg_icon_layers,  PURPLE,      "Reusable Skills\n& Accelerators",
     "25+ AON skills + 200+ enterprise Skills COE catalog — written once, enforced on every commit",
     "Zero knowledge drift; standards encoded Sprint 0 and carried unchanged to production",
     "15–20%", "additive gain from reusability", "Skills COE — 100,000+ users"),
    (svg_icon_gear,    DEEP_PURPLE, "Repetitive Task\nAutomation",
     "SonarQube, Fortify, OWASP, RBAC matrix, NFR matrix — automated per diff in 30–90 seconds",
     "Manual scan: 2–4 hrs per PR → diff-targeted scan: 30 sec; no full-project noise from legacy code",
     "22–32%", "effort saved (automation lever)", "LFG AI Readiness Q2 — 5 levers"),
    (svg_icon_shield,  PURPLE,      "AI Testing &\nTest Agents",
     "116 tests auto-generated; /run golden-path enforced; /migration-safe on every schema change",
     "AI code acceptance target 90%+; test suites generated from spec — quality loop closed before merge",
     "85%", "QA effort reduction", "U.S. Commercial Insurer — LFG RFI"),
    (svg_icon_refresh, DEEP_PURPLE, "React Front‑end\nPatterns",
     "7 role-scoped screens on Webpack 5 MFE shell — /i18n-parity, /wcag-audit, /uui-compliance live",
     "@aon/united-ui swap path documented in CLAUDE.md — production switch is a config change, not a rewrite",
     "20–30%", "design/build cycle reduction", "LFG AI Readiness North Star"),
]

ry_e = CY
for icon_fn, accent, drv_name, mech_title, mech_detail, pct, pct_label, source in eff_drivers:
    add_rect(s_eff, ML, ry_e, CW, row_h_e, fill=BG_LIGHT)
    add_rect(s_eff, ML, ry_e, 0.05, row_h_e, fill=accent)
    add_svg(s_eff, icon_fn(20, "#5F0095"),
            ML + 0.12, ry_e + (row_h_e - 0.24) / 2, 0.24, 0.24, 192)
    add_tb(s_eff, ML + 0.44, ry_e + 0.10, drv_w - 0.52, row_h_e - 0.18,
           drv_name, size=14, bold=True, color=accent)
    add_rect(s_eff, mech_x_e,       ry_e + 0.10, 0.02, row_h_e - 0.20, fill=LIGHT_GRAY)
    # title: generous h=0.50" so 2-line wrap renders cleanly
    add_tb(s_eff, mech_x_e + 0.10, ry_e + 0.06, mech_w - 0.16, 0.50,
           mech_title, size=14, bold=True, color=DEEP_PURPLE)
    add_tb(s_eff, mech_x_e + 0.10, ry_e + 0.58, mech_w - 0.16, row_h_e - 0.64,
           mech_detail, size=14, color=TEXT_BODY)
    add_rect(s_eff, imp_x_e,        ry_e + 0.10, 0.02, row_h_e - 0.20, fill=LIGHT_GRAY)
    add_tb(s_eff, imp_x_e + 0.12, ry_e + 0.06, imp_w - 0.16, 0.34,
           pct, size=20, bold=True, color=accent)
    add_tb(s_eff, imp_x_e + 0.12, ry_e + 0.44, imp_w - 0.16, 0.26,
           pct_label, size=14, color=TEXT_BODY)
    add_tb(s_eff, imp_x_e + 0.12, ry_e + 0.72, imp_w - 0.16, row_h_e - 0.78,
           source, size=14, italic=True, color=TEXT_SUB)
    ry_e += row_h_e + gap_eff

add_svg(s_eff, svg_icon_target(28, "#7E00FF"), ML + CW - 0.38, CY + 0.04, 0.30, 0.30, 256)
add_gt_to_slide(s_eff)
set_footer(s_eff)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12c: VALUE TRAJECTORY — 4-column sprint milestone
# ═══════════════════════════════════════════════════════════════════
s_val = prs.slides.add_slide(layout_content)
set_slide_title(s_val, "Value accumulates sprint by sprint — tracking to 40–55%", "Efficiency and Value Realization")
set_lead(s_val, "Sprint 0 encodes the standard. Each sprint wave activates additional drivers. Cumulative gain tracks Lincoln Financial's own 40–55% AD efficiency target by production readiness.")

# Reference band anchoring to Lincoln Financial's own benchmark
ref_h_v = 0.44
add_rect(s_val, ML, CY, CW, ref_h_v, fill=LP_BG)
add_rect(s_val, ML, CY, CW, 0.03, fill=PURPLE)
add_tb(s_val, ML + 0.16, CY + 0.09, CW - 0.28, 0.28,
       "Lincoln Financial AI Readiness (Q1/Q2): 40–55% AD gain over 24–36 months  ·  50–70% combined AD+AMS over 3–5 years",
       size=14, italic=True, color=DEEP_PURPLE)

n_ms    = 4
gap_ms  = 0.12
col_w_v = (CW - (n_ms - 1) * gap_ms) / n_ms   # 3.035"
col_y_v = CY + ref_h_v + 0.08                 # CY + 0.52"
note_h  = 0.38
col_h_v = AH - ref_h_v - 0.08 - 0.10 - note_h  # 4.20"

milestones = [
    ("Sprint 0",         "10–15%", DEEP_PURPLE,
     ["/init: CLAUDE.md loaded",
      "/security-review live",
      "51 NFRs + rules encoded",
      "Context from day one"],
     "Quality gate + context loop eliminated"),
    ("Sprint 2",         "25–30%", PURPLE,
     ["25+ custom skills live",
      "Explore + Plan concurrent",
      "Loop cut: 40 to 60%",
      "FNOL, scope, WCAG, perf"],
     "Parallel execution unlocked"),
    ("Sprint 4",         "35–40%", DEEP_PURPLE,
     ["RBAC + NFR matrix auto",
      "/run before task done",
      "116 tests per sprint",
      "/migration-safe on"],
     "85% QA effort reduction (comparable)"),
    ("Production Ready", "40–55%", PURPLE,
     ["5 swap paths checked",
      "Deploy gates wired",
      "POC to prod: checklist",
      "100,000+ scale live"],
     "Matches LFG AI Readiness AD target"),
]

# Column content vertical positions (relative to col_y_v)
hdr_h_v   = 0.54
pct_off   = 0.62   # % number y-offset from col top
pct_h_v   = 0.38
cg_off    = 1.08   # "cumulative gain" label
div1_off  = 1.36
act_off   = 1.46   # "What activates:" label
bul_start = 1.76   # first bullet
bul_step  = 0.24
div2_off  = bul_start + 4 * bul_step + 0.10   # 3.22"
eff_off   = div2_off + 0.14                    # 3.36"

for i, (ms_label, pct, col_col, bullets, effect) in enumerate(milestones):
    cx_v = ML + i * (col_w_v + gap_ms)
    # column background
    add_rect(s_val, cx_v, col_y_v, col_w_v, col_h_v, fill=BG_LIGHT)
    # colored header
    add_rect(s_val, cx_v, col_y_v, col_w_v, hdr_h_v, fill=col_col)
    add_tb(s_val, cx_v + 0.12, col_y_v + 0.12, col_w_v - 0.22, 0.30,
           ms_label, size=14, bold=True, color=WHITE)
    # % callout
    add_tb(s_val, cx_v + 0.12, col_y_v + pct_off, col_w_v - 0.22, pct_h_v,
           pct, size=22, bold=True, color=col_col)
    add_tb(s_val, cx_v + 0.12, col_y_v + cg_off, col_w_v - 0.22, 0.22,
           "cumulative gain", size=14, italic=True, color=TEXT_SUB)
    # divider 1
    add_rect(s_val, cx_v + 0.10, col_y_v + div1_off, col_w_v - 0.20, 0.02, fill=LIGHT_GRAY)
    # What activates label
    add_tb(s_val, cx_v + 0.12, col_y_v + act_off, col_w_v - 0.22, 0.26,
           "What activates:", size=14, bold=True, color=DEEP_PURPLE)
    # bullets
    for b_i, bul in enumerate(bullets):
        add_tb(s_val, cx_v + 0.12, col_y_v + bul_start + b_i * bul_step,
               col_w_v - 0.22, 0.22,
               f"- {bul}", size=14, color=TEXT_BODY)
    # divider 2
    add_rect(s_val, cx_v + 0.10, col_y_v + div2_off, col_w_v - 0.20, 0.02, fill=LIGHT_GRAY)
    # effect line
    add_tb(s_val, cx_v + 0.12, col_y_v + eff_off, col_w_v - 0.22, col_h_v - eff_off - 0.10,
           effect, size=14, bold=True, color=col_col)
    # connecting arrow
    if i < n_ms - 1:
        arr_x_v = cx_v + col_w_v + 0.01
        arr_y_v = col_y_v + col_h_v / 2 - 0.12
        add_tb(s_val, arr_x_v, arr_y_v, gap_ms - 0.02, 0.24,
               "▶", size=14, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# footnote
note_y_v = col_y_v + col_h_v + 0.10
add_tb(s_val, ML, note_y_v, CW, note_h,
       "※  Benchmarks sourced from Lincoln Financial AI Readiness RFI (Q1/Q2, July 2026) and Accenture case studies. "
       "POC-specific figures (116 tests, 45 endpoints, 25+ skills) are evidenced from the running codebase.",
       size=14, italic=True, color=TEXT_SUB)

add_svg(s_val, svg_icon_rocket(28, "#7E00FF"), ML + CW - 0.38, CY + 0.04, 0.30, 0.30, 256)
add_gt_to_slide(s_val)
set_footer(s_val)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13: CLOSE  — Pattern O (Bullet Memo)
# ═══════════════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(layout_content)
set_slide_title(s13, "We proved this model works — now we scale it", "What We Bring Forward")
set_lead(s13, "The demo proves the delivery model at component level. This engagement is where we prove it at production scale.")

close_bullets = [
    ("We shipped", "45 endpoints, 12 routers, and 116 tests — not with AI assistance as a claim, but as a demonstrated fact evidenced by a running codebase."),
    ("We encoded", "the security rules, architecture obligations, and production swap paths once in CLAUDE.md. Every engineer, every session, every sprint inherits them automatically."),
    ("We built", "skills that run SonarQube, Fortify, OWASP, and eight domain-specific gates on every diff — not as a one-time audit, but as the default delivery cadence."),
    ("We wired", "four specialist agents to run in parallel — cutting the research-then-build loop by 40–60% on cross-cutting changes."),
    ("We documented", "every non-obvious decision in ADRs, kept Architecture.md in sync with every change, and left five explicit production swap paths — so the path from POC to production is a checklist, not a rewrite."),
    ("We brought", "the Accenture × Anthropic enterprise Skills methodology — 200+ certified Skills in production, 100,000+ users deployed globally. This POC runs on the operating model already proven at enterprise scale."),
]

hl_h = 0.65
body_h = AH - hl_h - 0.15
tb_close = s13.shapes.add_textbox(Inches(ML), Inches(CY), Inches(CW), Inches(body_h))
tf_close = tb_close.text_frame
tf_close.word_wrap = True

for i, (bold_part, rest) in enumerate(close_bullets):
    p = tf_close.add_paragraph() if i > 0 else tf_close.paragraphs[0]
    p.space_before = Pt(18 if i > 0 else 0)
    p.line_spacing = Pt(18)
    r1 = p.add_run()
    r1.text = f"• {bold_part} "
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r1.font.name = FONT
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(14)
    r2.font.color.rgb = TEXT_BODY
    r2.font.name = FONT

hl_y = BY - hl_h
add_rect(s13, ML, hl_y, CW, hl_h, fill=DEEP_PURPLE)
add_tb(s13, ML + 0.20, hl_y + 0.15, CW - 0.40, hl_h - 0.18,
       "The path from POC to production is a checklist, not a rewrite — and the methodology behind it deploys at 100,000+ scale.",
       size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_gt_to_slide(s13)
set_footer(s13)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 16: APPENDIX — DEMO PERSONA MAP
# ═══════════════════════════════════════════════════════════════════
s_per = prs.slides.add_slide(layout_content)
set_slide_title(s_per, "Seven personas — roles, scope, and entitlements", "Appendix: Demo Persona Map")
set_lead(s_per, "Each persona maps to a named RFP user type. Org scope is enforced server-side via JWT — the persona switch produces a visibly different experience, not just a label change.")

persona_data = [
    # (name, role, level_label, org_node, accent_color, description)
    ("Sarah Whitfield",  "C-Suite",                    "Corporate",  "CORP-HOSP",     DEEP_PURPLE,
     "Full portfolio visibility: all claims, PII, analytics, export, restricted data"),
    ("Daniel Osei",      "Risk Manager · Client Admin","Corporate",  "CORP-HOSP",     DEEP_PURPLE,
     "Full access + FNOL submission, document upload, tenant administration"),
    ("Priya Raman",      "Location Manager · Airport", "Location",   "LOC-JFK",       PURPLE,
     "Location-scoped: claims viewer, analytics, FNOL, document management"),
    ("Marcus Lindqvist", "Functional Lead · HRBP",     "Location",   "LOC-JFK",       PURPLE,
     "Location-scoped: claims viewer and FNOL submission only"),
    ("Maria Santos",     "Site Manager · Restaurant",  "Site",       "SITE-JFK-T4",   LIGHT_GRAY,
     "Site-scoped: claims viewer, FNOL submission, document upload"),
    ("Tom Beckett",      "Employee · Reporter",        "Site",       "SITE-JFK-T4",   LIGHT_GRAY,
     "FNOL submission and own-claims view only (claims_own_only enforced)"),
    ("Unassigned User",  "Unauthorised User",          "No Access",  "—",             BG_LIGHT,
     "No org node assigned — BR-001 returns 403 on all data routes"),
]

col_w_ps  = (CW - 0.15) / 2          # 6.175"
n_left    = 4
n_right   = len(persona_data) - n_left  # 3
gap_ps    = 0.08
cell_h_ps = (AH - (n_left - 1) * gap_ps) / n_left  # 1.24"

right_total_h = n_right * cell_h_ps + (n_right - 1) * gap_ps
right_start_y = CY + (AH - right_total_h) / 2

badge_w_ps = 1.55

for i, (name, role, level, org_node, accent, desc) in enumerate(persona_data):
    if i < n_left:
        cx      = ML
        cy_card = CY + i * (cell_h_ps + gap_ps)
    else:
        cx      = ML + col_w_ps + 0.15
        cy_card = right_start_y + (i - n_left) * (cell_h_ps + gap_ps)

    add_rect(s_per, cx, cy_card, col_w_ps, cell_h_ps, fill=BG_LIGHT, border=LIGHT_GRAY)
    add_rect(s_per, cx, cy_card, 0.07, cell_h_ps, fill=accent)
    # Name
    name_w = col_w_ps - badge_w_ps - 0.36
    add_tb(s_per, cx + 0.15, cy_card + 0.10, name_w, 0.25,
           name, size=14, bold=True, color=DEEP_PURPLE)
    # Role
    add_tb(s_per, cx + 0.15, cy_card + 0.37, col_w_ps - 0.22, 0.22,
           role, size=14, italic=True, color=TEXT_SUB)
    # Description
    add_tb(s_per, cx + 0.15, cy_card + 0.62, col_w_ps - 0.22, cell_h_ps - 0.72,
           desc, size=14, color=TEXT_BODY)
    # Org-level badge (top-right)
    bx_ps = cx + col_w_ps - badge_w_ps - 0.08
    add_rect(s_per, bx_ps, cy_card + 0.10, badge_w_ps, 0.26, fill=LP_BG)
    add_tb(s_per, bx_ps + 0.06, cy_card + 0.10, badge_w_ps - 0.08, 0.26,
           f"{level}  ·  {org_node}", size=14, color=PURPLE)

add_svg(s_per, svg_icon_network(32, "#7E00FF"), ML + CW - 0.45, CY + 0.06, 0.36, 0.36, 256)
add_gt_to_slide(s_per)
set_footer(s_per)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 17: ORG HIERARCHY + RBAC MODEL — 3-column
# ═══════════════════════════════════════════════════════════════════
s_rbac = prs.slides.add_slide(layout_content)
set_slide_title(s_rbac, "Scope × capability — how RBAC governs every request", "Access Architecture")
set_lead(s_rbac, "The JWT determines where you sit in the client hierarchy. Your groups determine what you can do. BR-001 enforces the intersection server-side on every API call.")

# 2-column layout: Left = Org Hierarchy (personas embedded), Right = RBAC Rules
left_w_rb  = 4.80
gap_mid_rb = 0.20
right_w_rb = CW - left_w_rb - gap_mid_rb   # 7.50"
right_x_rb = ML + left_w_rb + gap_mid_rb   # 5.42"

# ── Left col: Client Org Hierarchy ───────────────────────────────────
add_rect(s_rbac, ML, CY, left_w_rb, AH, fill=WHITE, border=LIGHT_GRAY)
add_rect(s_rbac, ML, CY, left_w_rb, 0.05, fill=DEEP_PURPLE)
add_svg(s_rbac, svg_icon_network(24, "#7E00FF"), ML + 0.10, CY + 0.12, 0.26, 0.26, 256)
add_tb(s_rbac, ML + 0.44, CY + 0.14, left_w_rb - 0.54, 0.26,
       "Client Org Hierarchy", size=14, bold=True, color=PURPLE)

lhdr_h   = 0.50
larr_h   = 0.25
avail_lb = AH - lhdr_h - 2 * larr_h   # 4.20"
lbox_h   = avail_lb / 3               # 1.40"

org_levels = [
    ("Corporate",  "CORP-HOSP",    "Hospitality Ltd",   DEEP_PURPLE, WHITE,
     "Sarah Whitfield  ·  Daniel Osei",   "Full portfolio — all locations and sites"),
    ("Location",   "LOC-JFK",      "JFK Airport",       PURPLE,      WHITE,
     "Priya Raman  ·  Marcus Lindqvist",  "Scoped to own location and its sites"),
    ("Site",       "SITE-JFK-T4",  "Bistro Restaurant", LP_BG,       DEEP_PURPLE,
     "Maria Santos  ·  Tom Beckett",      "Scoped to own site claims only"),
]

bx_rb = ML + 0.10
bw_rb = left_w_rb - 0.20
y_lev = CY + lhdr_h
for j, (level, node, org, bg, tc, personas, scope) in enumerate(org_levels):
    add_rect(s_rbac, bx_rb, y_lev, bw_rb, lbox_h, fill=bg)
    add_tb(s_rbac, bx_rb + 0.12, y_lev + 0.10, bw_rb - 0.22, 0.26,
           f"{level}  ·  {node}", size=14, bold=True, color=tc)
    add_tb(s_rbac, bx_rb + 0.12, y_lev + 0.38, bw_rb - 0.22, 0.24,
           org, size=14, italic=True, color=tc)
    add_tb(s_rbac, bx_rb + 0.12, y_lev + 0.64, bw_rb - 0.22, 0.26,
           personas, size=14, color=tc)
    add_tb(s_rbac, bx_rb + 0.12, y_lev + 0.94, bw_rb - 0.22, lbox_h - 1.02,
           scope, size=14, bold=True, color=tc)
    y_lev += lbox_h
    if j < 2:
        add_tb(s_rbac, bx_rb, y_lev, bw_rb, larr_h,
               "▼  scope inherits downward", size=14, italic=True,
               color=PURPLE, align=PP_ALIGN.CENTER)
        y_lev += larr_h

# ── Right col: BR-001 Enforcement ────────────────────────────────────
add_rect(s_rbac, right_x_rb, CY, right_w_rb, AH, fill=WHITE, border=LIGHT_GRAY)
add_rect(s_rbac, right_x_rb, CY, right_w_rb, 0.05, fill=DEEP_PURPLE)
add_svg(s_rbac, svg_icon_shield(24, "#7E00FF"), right_x_rb + 0.10, CY + 0.12, 0.26, 0.26, 256)
add_tb(s_rbac, right_x_rb + 0.44, CY + 0.14, right_w_rb - 0.54, 0.26,
       "BR-001 Enforcement — 4 rules applied to every API request", size=14, bold=True, color=PURPLE)

rhdr_h  = 0.50
rnote_h = 0.38
rgap    = 0.10
avail_r = AH - rhdr_h - 0.04 - rnote_h - 3 * rgap   # 3.98"
rule_h2 = avail_r / 4                                  # 0.995"

rbac_rules = [
    ("01", "JWT carries org_node",
     "Injected at login — never read from query params, headers, or request body"),
    ("02", "current_scope dependency",
     "FastAPI dep on every data route — routes without it are automatically flagged unprotected"),
    ("03", "SQL materialized-path filter",
     "WHERE org_node LIKE '{user_node}%' — child nodes fall in scope automatically"),
    ("04", "Out-of-scope = HTTP 403, not 404",
     "404 leaks existence; 403 denies without disclosure — audit log written on every deny"),
]
ry      = CY + rhdr_h + 0.04
badge_w = 0.44
text_rx = right_x_rb + 0.10 + badge_w + 0.10
text_rw = right_w_rb - 0.10 - badge_w - 0.10 - 0.10   # 6.76"

for num, rtitle, rdesc in rbac_rules:
    add_rect(s_rbac, right_x_rb + 0.10, ry, right_w_rb - 0.20, rule_h2, fill=BG_LIGHT)
    add_rect(s_rbac, right_x_rb + 0.10, ry, badge_w, rule_h2, fill=DEEP_PURPLE)
    add_tb(s_rbac, right_x_rb + 0.10, ry + (rule_h2 - 0.28) / 2, badge_w, 0.28,
           num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(s_rbac, text_rx, ry + 0.10, text_rw, 0.28,
           rtitle, size=14, bold=True, color=DEEP_PURPLE)
    add_tb(s_rbac, text_rx, ry + 0.42, text_rw, rule_h2 - 0.52,
           rdesc, size=14, color=TEXT_BODY)
    ry += rule_h2 + rgap

add_tb(s_rbac, right_x_rb + 0.10, ry, right_w_rb - 0.20, CY + AH - ry,
       "※  groups_csv determines what each user can DO within their org scope — 10 capability groups on the next slide",
       size=14, italic=True, color=TEXT_SUB)

add_gt_to_slide(s_rbac)
set_footer(s_rbac)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 17b: RBAC OWNERSHIP — 3 fixed constraints + 4 owners
# ═══════════════════════════════════════════════════════════════════
s_own = prs.slides.add_slide(layout_content)
set_slide_title(s_own, "Three constraints, four owners — all RBAC decisions", "Access Architecture")
set_lead(s_own, "NFR-33, BR-001 and MDM/Reltio are load-bearing — the RBAC build inherits these constraints, not the other way around.")

own_lw = 4.80; own_gap = 0.20
own_rw = CW - own_lw - own_gap        # 7.50"
own_rx = ML + own_lw + own_gap        # 5.42"

# ── Left: 3 fixed constraints ──────────────────────────────────────
add_rect(s_own, ML, CY, own_lw, AH, fill=WHITE, border=LIGHT_GRAY)
add_rect(s_own, ML, CY, own_lw, 0.05, fill=DEEP_PURPLE)
add_svg(s_own, svg_icon_shield(22, "#7E00FF"), ML + 0.10, CY + 0.12, 0.24, 0.24, 192)
add_tb(s_own, ML + 0.42, CY + 0.14, own_lw - 0.52, 0.26,
       "Three constraints that cannot change", size=14, bold=True, color=PURPLE)

constraints_rb = [
    ("01", "Okta is the only IdP",
     "NFR-33 and the Meridian Pattern prohibit custom identity stores.\nEvery user — internal and external — authenticates through Okta."),
    ("02", "JWT carries exactly two RBAC claims",
     "groups — the role assigned. org_node — the organisational scope,\nsourced from MDM/Reltio. App reads only the token, never the client."),
    ("03", "Scope inheritance is strictly downward",
     "BR-001: user's assigned node plus all descendants.\nPeers and parents excluded. Enforced at Claims Copilot on every API call."),
]

n_crb = len(constraints_rb); gap_crb = 0.12
card_crb_h = (AH - 0.44 - (n_crb - 1) * gap_crb) / n_crb    # ≈ 1.50"
ycrb = CY + 0.44
for num_crb, ctitle_crb, cdesc_crb in constraints_rb:
    add_rect(s_own, ML + 0.08, ycrb, own_lw - 0.16, card_crb_h, fill=LP_BG)
    add_rect(s_own, ML + 0.08, ycrb, 0.40, card_crb_h, fill=DEEP_PURPLE)
    add_tb(s_own, ML + 0.08, ycrb + (card_crb_h - 0.28) / 2, 0.40, 0.28,
           num_crb, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(s_own, ML + 0.56, ycrb + 0.10, own_lw - 0.72, 0.30,
           ctitle_crb, size=14, bold=True, color=DEEP_PURPLE)
    add_tb(s_own, ML + 0.56, ycrb + 0.44, own_lw - 0.72, card_crb_h - 0.54,
           cdesc_crb, size=14, color=TEXT_BODY)
    ycrb += card_crb_h + gap_crb

# ── Right: ownership table ─────────────────────────────────────────
add_rect(s_own, own_rx, CY, own_rw, AH, fill=WHITE, border=LIGHT_GRAY)
add_rect(s_own, own_rx, CY, own_rw, 0.05, fill=DEEP_PURPLE)
add_svg(s_own, svg_icon_layers(22, "#7E00FF"), own_rx + 0.10, CY + 0.12, 0.24, 0.24, 192)
add_tb(s_own, own_rx + 0.42, CY + 0.14, own_rw - 0.52, 0.26,
       "Four ownership layers — who controls what", size=14, bold=True, color=PURPLE)

c1w_o, c2w_o, c3w_o = 1.80, 2.40, 3.10    # total = 7.30 = own_rw - 0.20
c1x_o = own_rx + 0.10
c2x_o = c1x_o + c1w_o
c3x_o = c2x_o + c2w_o

add_rect(s_own, c1x_o, CY + 0.44, own_rw - 0.20, 0.36, fill=DEEP_PURPLE)
for lbl_o, cx_o in [("Layer", c1x_o + 0.08), ("Owner", c2x_o + 0.08), ("Controls", c3x_o + 0.08)]:
    add_tb(s_own, cx_o, CY + 0.50, 1.60, 0.24, lbl_o, size=14, bold=True, color=WHITE)

owners_rb = [
    ("Platform Identity",     "Aon IT / Okta admin",         "Okta groups, SCIM provisioning, MFA enforcement, token lifecycle, JWKS endpoint"),
    ("Role & Privilege Defn", "App Owner + Aon Security",    "Which roles exist, privilege matrix, document audience classification, change sign-off"),
    ("Client-Side Access",    "Persona 2 — Client Admin",    "Who in their org gets which role; provisioning and updates within their authorised scope only"),
    ("Aon Provisioning",      "Aon Claims Professional",     "Creates external users in Meridian; triggers Okta registration email and MFA setup (NFR-03)"),
]

own_body_rb = AH - 0.44 - 0.36
row_own_h = own_body_rb / len(owners_rb)    # ≈ 1.10"
yo_rb = CY + 0.80
for i_o, (layer_o, owner_o, controls_o) in enumerate(owners_rb):
    fill_o = BG_LIGHT if i_o % 2 == 0 else WHITE
    add_rect(s_own, c1x_o, yo_rb, own_rw - 0.20, row_own_h, fill=fill_o)
    add_tb(s_own, c1x_o + 0.08, yo_rb + 0.10, c1w_o - 0.14, row_own_h - 0.16,
           layer_o, size=14, bold=True, color=DEEP_PURPLE)
    add_rect(s_own, c2x_o, yo_rb, 0.02, row_own_h, fill=LIGHT_GRAY)
    add_tb(s_own, c2x_o + 0.08, yo_rb + 0.10, c2w_o - 0.14, row_own_h - 0.16,
           owner_o, size=14, color=TEXT_BODY)
    add_rect(s_own, c3x_o, yo_rb, 0.02, row_own_h, fill=LIGHT_GRAY)
    add_tb(s_own, c3x_o + 0.08, yo_rb + 0.10, c3w_o - 0.14, row_own_h - 0.16,
           controls_o, size=14, color=TEXT_BODY)
    yo_rb += row_own_h

add_gt_to_slide(s_own)
set_footer(s_own)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 17c: ACCESS JOURNEY — 5 steps to a live JWT claim
# ═══════════════════════════════════════════════════════════════════
s_journey = prs.slides.add_slide(layout_content)
set_slide_title(s_journey, "Five steps from access request to live JWT claim", "Access Architecture")
set_lead(s_journey, "External client users enter through Meridian — provisioned by an internal Aon user, authenticated via Okta PKCE, and scoped by the MDM hierarchy on every API call.")

access_steps = [
    (DEEP_PURPLE, "Pre-condition — MDM hierarchy confirmed in Reltio",
     "Corporate → Region → Location loaded; all org_node values established before any user can be scoped. R-14 (MDM → Okta binding mechanism and refresh frequency) must be resolved at mobilisation workshops before this step is complete."),
    (PURPLE, "Client Admin (Persona 2) raises the access request",
     "Specifies role (one of 7 personas), the org node to scope the user to, and privilege groups required: eFNOL, claim list visibility, documents, contacts, or insights."),
    (DEEP_PURPLE, "Aon Claims Professional provisions the user in Meridian",
     "Creates the external user account in Meridian; Okta sends a registration email; user sets up MFA and Okta credentials on first login. Internal users hold the provisioning right per NFR-03."),
    (PURPLE, "Okta group assigned via SCIM — org_node claim bound",
     "User added to the correct Okta group (one group per role). The org_node custom claim populated from the MDM hierarchy binding. OI-26 governs the exact Okta group structure and custom claim configuration."),
    (DEEP_PURPLE, "Authenticate via PKCE → JWT issued → role-scoped screen activates",
     "Meridian SPA initiates Okta PKCE flow; JWT carries groups + org_node. BR-001 scope enforcement fires on every API route. Correct 1-of-7 screen renders. Deep links preserved via OAuth2 state parameter (F-MER-04)."),
]

n_as = len(access_steps); gap_as = 0.08
step_as_h = (AH - (n_as - 1) * gap_as) / n_as    # ≈ 0.976"
badge_as_w = 0.44
text_as_x = ML + badge_as_w + 0.10
text_as_w = CW - badge_as_w - 0.18

yas = CY
for i_as, (accent_as, stitle_as, sdetail_as) in enumerate(access_steps):
    add_rect(s_journey, ML, yas, CW, step_as_h, fill=BG_LIGHT)
    add_rect(s_journey, ML, yas, badge_as_w, step_as_h, fill=accent_as)
    add_tb(s_journey, ML, yas + (step_as_h - 0.28) / 2, badge_as_w, 0.28,
           f"0{i_as + 1}", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(s_journey, text_as_x, yas + 0.08, text_as_w, 0.30,
           stitle_as, size=14, bold=True, color=accent_as)
    add_tb(s_journey, text_as_x, yas + 0.42, text_as_w, step_as_h - 0.50,
           sdetail_as, size=14, color=TEXT_BODY)
    yas += step_as_h + gap_as

add_svg(s_journey, svg_icon_network(26, "#7E00FF"),
        ML + CW - 0.38, CY + 0.06, 0.32, 0.32, 256)
add_gt_to_slide(s_journey)
set_footer(s_journey)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 17d: GOVERNANCE + DOCUMENT RBAC + OPEN BLOCKERS
# ═══════════════════════════════════════════════════════════════════
s_gov2 = prs.slides.add_slide(layout_content)
set_slide_title(s_gov2, "Governance, document RBAC, and the two open blockers", "Access Architecture")
set_lead(s_gov2, "Document access runs on two enforcement layers below role scope. Two open items must close at mobilisation before the RBAC design deliverable can be baselined.")

top_h_g = 3.46
blocker_gap_g = 0.10
blocker_h_g = AH - top_h_g - blocker_gap_g    # ≈ 1.64"

doc_lw_g = 5.80; gov_gap_g = 0.20
gov_rw_g = CW - doc_lw_g - gov_gap_g          # 6.50"
gov_rx_g = ML + doc_lw_g + gov_gap_g          # 6.62"

# ── Left: Document RBAC ────────────────────────────────────────────
add_rect(s_gov2, ML, CY, doc_lw_g, top_h_g, fill=WHITE, border=LIGHT_GRAY)
add_rect(s_gov2, ML, CY, doc_lw_g, 0.05, fill=DEEP_PURPLE)
add_svg(s_gov2, svg_icon_shield(22, "#7E00FF"), ML + 0.10, CY + 0.12, 0.24, 0.24, 192)
add_tb(s_gov2, ML + 0.42, CY + 0.14, doc_lw_g - 0.52, 0.26,
       "Document RBAC — two enforcement layers", size=14, bold=True, color=PURPLE)

doc_lvl_h = (top_h_g - 0.44 - 0.06) / 2    # ≈ 1.48" each

lv1_y = CY + 0.44
add_rect(s_gov2, ML + 0.08, lv1_y, doc_lw_g - 0.16, doc_lvl_h, fill=LP_BG)
add_tb(s_gov2, ML + 0.16, lv1_y + 0.08, doc_lw_g - 0.30, 0.28,
       "Level 1 — Org-node scope (BR-001)", size=14, bold=True, color=DEEP_PURPLE)
add_tb(s_gov2, ML + 0.16, lv1_y + 0.40, doc_lw_g - 0.30, doc_lvl_h - 0.48,
       "User sees only documents on claims within their org_node. Same downward scope logic as claims and analytics. S-DMS enforces at the metadata layer before returning any ECM reference.",
       size=14, color=TEXT_BODY)

lv2_y = lv1_y + doc_lvl_h + 0.06
add_rect(s_gov2, ML + 0.08, lv2_y, doc_lw_g - 0.16, doc_lvl_h, fill=BG_LIGHT)
add_tb(s_gov2, ML + 0.16, lv2_y + 0.08, doc_lw_g - 0.30, 0.28,
       "Level 2 — Document attribute enforcement (S-DMS)", size=14, bold=True, color=DEEP_PURPLE)
add_tb(s_gov2, ML + 0.16, lv2_y + 0.40, doc_lw_g - 0.30, doc_lvl_h - 0.48,
       "Attributes: Default, Internal, View on Web, Modify on Web, Access Controlled. Audience: Client-visible, Internal, Carrier-only. BR-007/008: provenance tag blocks client docs from Aon Brokers.",
       size=14, color=TEXT_BODY)

# ── Right: Ongoing governance ──────────────────────────────────────
add_rect(s_gov2, gov_rx_g, CY, gov_rw_g, top_h_g, fill=WHITE, border=LIGHT_GRAY)
add_rect(s_gov2, gov_rx_g, CY, gov_rw_g, 0.05, fill=DEEP_PURPLE)
add_svg(s_gov2, svg_icon_gear(22, "#7E00FF"), gov_rx_g + 0.10, CY + 0.12, 0.24, 0.24, 192)
add_tb(s_gov2, gov_rx_g + 0.42, CY + 0.14, gov_rw_g - 0.52, 0.26,
       "Three ongoing governance requirements", size=14, bold=True, color=PURPLE)

gov_items_g = [
    ("Quarterly access recertification",
     "Okta per-group reports quarterly. Persona 2 certifies client users; Claims Ops certifies internal provisioning."),
    ("Privilege matrix as a living document",
     "Role permission changes require App Owner + Aon Security sign-off. M2 privilege matrix is the live source of truth."),
    ("Penetration test gate — Dec 2026",
     "Aon executes Dec 2026. 7 persona accounts required. Negative scope tests for BR-001/003/007/008 mandatory."),
]

n_gov_g = len(gov_items_g); gap_gov_g = 0.10
row_gov_h = (top_h_g - 0.44 - (n_gov_g - 1) * gap_gov_g) / n_gov_g    # ≈ 0.94"
yg = CY + 0.44
for i_g, (gtitle_g, gdesc_g) in enumerate(gov_items_g):
    fill_g = LP_BG if i_g % 2 == 0 else BG_LIGHT
    add_rect(s_gov2, gov_rx_g + 0.08, yg, gov_rw_g - 0.16, row_gov_h, fill=fill_g)
    add_rect(s_gov2, gov_rx_g + 0.08, yg, 0.04, row_gov_h, fill=PURPLE)
    add_tb(s_gov2, gov_rx_g + 0.20, yg + 0.08, gov_rw_g - 0.36, 0.28,
           gtitle_g, size=14, bold=True, color=DEEP_PURPLE)
    add_tb(s_gov2, gov_rx_g + 0.20, yg + 0.40, gov_rw_g - 0.36, row_gov_h - 0.48,
           gdesc_g, size=14, color=TEXT_BODY)
    yg += row_gov_h + gap_gov_g

# ── Bottom: 2 open blockers ────────────────────────────────────────
blocker_y_g = CY + top_h_g + blocker_gap_g    # 5.21"
bw_g = (CW - 0.20) / 2                         # 6.15" each
b2x_g = ML + bw_g + 0.20                       # 6.77"

blockers_g = [
    ("R-14", "MDM → Okta org scope binding and refresh frequency",
     "Resolve at mobilisation. Determines entitlement latency and the joiner / mover / leaver revocation window — how quickly access changes propagate when a user changes role, moves node, or leaves."),
    ("OI‑26", "Okta configuration for Client User groups and claims",
     "Resolve at mobilisation. Determines RBAC build effort and client onboarding process design. Defines exact Okta group structure and org_node custom claim implementation."),
]

for k_g, (ref_g, btitle_g, bdesc_g) in enumerate(blockers_g):
    bx_g = ML if k_g == 0 else b2x_g
    add_rect(s_gov2, bx_g, blocker_y_g, bw_g, blocker_h_g, fill=LP_BG, border=PURPLE)
    add_rect(s_gov2, bx_g, blocker_y_g, 0.05, blocker_h_g, fill=PURPLE)
    add_tb(s_gov2, bx_g + 0.14, blocker_y_g + 0.10, 0.90, 0.28,
           ref_g, size=14, bold=True, color=PURPLE)
    add_tb(s_gov2, bx_g + 0.14, blocker_y_g + 0.42, bw_g - 0.24, 0.28,
           btitle_g, size=14, bold=True, color=DEEP_PURPLE)
    add_tb(s_gov2, bx_g + 0.14, blocker_y_g + 0.74, bw_g - 0.24, blocker_h_g - 0.84,
           bdesc_g, size=14, color=TEXT_BODY)

add_gt_to_slide(s_gov2)
set_footer(s_gov2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE: Python FastAPI → .NET Core production delta
# ═══════════════════════════════════════════════════════════════════
s_net = prs.slides.add_slide(layout_content)
set_slide_title(s_net, "Python FastAPI → .NET Core: the key production delta", "Access Architecture")
set_lead(s_net, "The RFP mandates .NET Core as the Aon Enterprise Architecture API standard. Python FastAPI proves the integration model in the POC — production requires a .NET Core rewrite on Azure.")

n_net = 3; gap_net = 0.14
col_w_net = (CW - (n_net - 1) * gap_net) / n_net    # ≈ 4.073"

net_col_cfg = [
    (DEEP_PURPLE, svg_icon_shield, "The RFP requirement",
     [("Why it is required",
       "API-first, loosely coupled, versioned in .NET Core — possibly containerised. Verbatim requirement from the RFP, Must priority."),
      ("Why it already exists",
       "Claims Copilot API layer runs on .NET Core consistent with the Aon Enterprise Architecture standard stack."),
      ("Azure tooling standard",
       "Azure API Management + App Service are the approved hosting stack. Azure DevOps is mandatory for CI/CD. .NET Core fits natively."),
     ]),
    (PURPLE, svg_icon_code, "What the POC built",
     [("Runtime + API surface",
       "Python 3.12 + FastAPI + Pydantic v2. 45 endpoints · 12 routers · 116 passing tests. JWT HS256, BR-001 enforced on every route."),
      ("Role in delivery",
       "Proves the integration model and RBAC logic in the demo environment — the validated stand-in layer."),
      ("Production gap",
       "Python FastAPI is not in the Aon Enterprise Architecture approved stack. Azure deployment and RS256 JWT not yet configured."),
     ]),
    (DEEP_PURPLE, svg_icon_layers, "What production requires",
     [("Runtime + gateway",
       ".NET Core (C#), versioned and containerised. Azure API Management in front for rate limiting, auth, and versioning."),
      ("Security upgrade",
       "RS256 JWT validated against Okta JWKS endpoint. Same BR-001 scope enforcement — different implementation language."),
      ("Migration scope",
       "API contracts and RBAC rules are language-agnostic. This is an implementation change, not an architectural redesign. Sprint 5 to 6."),
     ]),
]

n_i_net = 3; gap_i_net = 0.14
item_h_net = (AH - 0.44 - 0.10 - (n_i_net - 1) * gap_i_net - 0.10) / n_i_net    # ≈ 1.427"

for ci_n, (accent_n, icon_fn_n, hdr_n, items_n) in enumerate(net_col_cfg):
    cx_n = ML + ci_n * (col_w_net + gap_net)
    add_rect(s_net, cx_n, CY, col_w_net, AH, fill=BG_LIGHT if ci_n != 1 else WHITE)
    add_rect(s_net, cx_n, CY, col_w_net, 0.05, fill=accent_n)
    add_svg(s_net, icon_fn_n(20, "#7E00FF"), cx_n + 0.10, CY + 0.12, 0.22, 0.22, 192)
    add_tb(s_net, cx_n + 0.40, CY + 0.14, col_w_net - 0.50, 0.26,
           hdr_n, size=14, bold=True, color=PURPLE)
    iy_n = CY + 0.44 + 0.10
    for i_n, (lbl_n, det_n) in enumerate(items_n):
        add_rect(s_net, cx_n + 0.08, iy_n, col_w_net - 0.16, item_h_net,
                 fill=LP_BG if ci_n != 1 else BG_LIGHT)
        add_rect(s_net, cx_n + 0.08, iy_n, 0.04, item_h_net, fill=accent_n)
        add_tb(s_net, cx_n + 0.20, iy_n + 0.08, col_w_net - 0.30, 0.28,
               lbl_n, size=14, bold=True, color=DEEP_PURPLE)
        add_tb(s_net, cx_n + 0.20, iy_n + 0.40, col_w_net - 0.30, item_h_net - 0.48,
               det_n, size=14, color=TEXT_BODY)
        iy_n += item_h_net + gap_i_net

add_gt_to_slide(s_net)
set_footer(s_net)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 18: CAPABILITY GROUPS REFERENCE — 2×5 card grid
# ═══════════════════════════════════════════════════════════════════
s_grp = prs.slides.add_slide(layout_content)
set_slide_title(s_grp, "Ten groups define what each persona can see and do", "Capability Groups")
set_lead(s_grp, "Every persona carries a groups_csv. Each group unlocks specific screens and actions — anything not granted is implicitly denied by the API.")

group_cards = [
    # LEFT col — view/access (DEEP_PURPLE accent)
    ("claims_viewer",
     "Browse claims list and open claim detail",
     "Scoped to org_node from JWT — BR-001 filter on every query",        DEEP_PURPLE),
    ("claims_analytics",
     "KPI dashboard: 17 tiles, trend charts, product breakdown",
     "Org-scoped aggregate only — no raw PII or individual records",       DEEP_PURPLE),
    ("claims_docs",
     "View documents after 3-gate Pillar 1 check",
     "Gates: audience · security_attr · provenance — ecm_ref never in response", DEEP_PURPLE),
    ("claims_view_pii",
     "Unmask: named_insured, loss_address, submitted_by",
     "Without this: PII fields display as [REDACTED] in all views",        DEEP_PURPLE),
    ("claims_view_restricted",
     "Access restricted-classification claim rows",
     "Without this: restricted rows hidden from list and detail entirely",  DEEP_PURPLE),
    # RIGHT col — action/modifier (PURPLE accent)
    ("claims_fnol",
     "Submit FNOL via wizard; manage cross-device drafts",
     "Outbox-first pattern — claim persists through backend outage",        PURPLE),
    ("claims_export",
     "Download filtered claims as CSV report",
     "PII fields masked in export without claims_view_pii",                 PURPLE),
    ("claims_upload_docs",
     "Attach documents to a claim record",
     "Org-scoped: cannot attach to out-of-scope claims",                    PURPLE),
    ("claims_client_admin",
     "Configure field registry (73 fields) and tenant settings",
     "Corporate-level org node required",                                   PURPLE),
    ("claims_own_only",
     "Modifier: restricts claims_viewer to own submitted claims",
     "Overrides org scope — only own FNOL history visible",                 PURPLE),
]

col_w_gr = (CW - 0.15) / 2   # 6.175"
n_rows_gr = 5
gap_gr    = 0.07
row_h_gr  = (AH - (n_rows_gr - 1) * gap_gr) / n_rows_gr   # 0.984"

for idx, (gname, desc, restriction, accent) in enumerate(group_cards):
    col_idx = idx // n_rows_gr
    row_idx = idx % n_rows_gr
    gx = ML + col_idx * (col_w_gr + 0.15)
    gy = CY + row_idx * (row_h_gr + gap_gr)
    add_rect(s_grp, gx, gy, col_w_gr, row_h_gr, fill=BG_LIGHT, border=LIGHT_GRAY)
    add_rect(s_grp, gx, gy, 0.07, row_h_gr, fill=accent)
    add_tb(s_grp, gx + 0.14, gy + 0.08, col_w_gr - 0.22, 0.24,
           gname, size=14, bold=True, color=DEEP_PURPLE)
    add_tb(s_grp, gx + 0.14, gy + 0.35, col_w_gr - 0.22, 0.28,
           desc, size=14, color=TEXT_BODY)
    add_tb(s_grp, gx + 0.14, gy + 0.64, col_w_gr - 0.22, row_h_gr - 0.72,
           restriction, size=14, italic=True, color=TEXT_SUB)

add_svg(s_grp, svg_icon_shield(28, "#7E00FF"), ML + CW - 0.44, CY + 0.06, 0.34, 0.34, 256)
add_gt_to_slide(s_grp)
set_footer(s_grp)


# ═══════════════════════════════════════════════════════════════════
# SAVE + VERIFY
# ═══════════════════════════════════════════════════════════════════
output_path = os.path.join(os.getcwd(), "AON_Claude_Code_Skills_Deck.pptx")
prs.save(output_path)
cleanup_temp()
print(f"Saved: {output_path}")

import subprocess
_venv_py = os.path.join(os.path.dirname(SKILL_DIR), ".venv", "Scripts", "python.exe")
_py = _venv_py if os.path.exists(_venv_py) else sys.executable

result = subprocess.run(
    [_py, os.path.join(SKILL_DIR, "scripts", "verify_pptx.py"), output_path],
    capture_output=True, text=True
)
print(result.stdout)
if result.returncode != 0:
    print("VERIFY ERRORS:", result.stderr)

thumb_dir = os.path.join(os.getcwd(), "thumbnails")
os.makedirs(thumb_dir, exist_ok=True)
thumb_result = subprocess.run(
    [_py, os.path.join(SKILL_DIR, "scripts", "thumbnail.py"), output_path, thumb_dir],
    capture_output=True, text=True
)
print(thumb_result.stdout)
if thumb_result.returncode != 0:
    print("THUMBNAIL ERRORS:", thumb_result.stderr)
